#!/usr/bin/env python3
"""Generate U1 Query Flow PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ─── Colors ──────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x0A, 0x29, 0x55)
BLUE    = RGBColor(0x17, 0x5C, 0xA6)
TEAL    = RGBColor(0x00, 0x7E, 0x8A)
GREEN   = RGBColor(0x1B, 0x6F, 0x42)
ORANGE  = RGBColor(0xC4, 0x50, 0x08)
PURPLE  = RGBColor(0x5A, 0x23, 0x8C)
DGRAY   = RGBColor(0x44, 0x44, 0x44)
LGRAY   = RGBColor(0xF0, 0xF2, 0xF5)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
L_ORG   = RGBColor(0xFF, 0xF0, 0xE8)
L_PUR   = RGBColor(0xF2, 0xEB, 0xFF)
L_TEAL  = RGBColor(0xE0, 0xF5, 0xF7)
QARROW  = BLUE
RARROW  = GREEN


# ─── Helpers ─────────────────────────────────────────────────────────────────

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rect(slide, x, y, w, h, fill, border=None, bw=0.5):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp


def box(slide, x, y, w, h, fill, lines, fsize=11, fc=WHITE, bold0=True):
    shp = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shp.line.width = Pt(0.5)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(fsize)
        r.font.color.rgb = fc
        r.font.bold = (bold0 and i == 0)
    return shp


def lbl(slide, x, y, w, h, text, fsize=9, fc=DGRAY, bold=False, italic=False,
        align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fsize)
    r.font.color.rgb = fc
    r.font.bold = bold
    r.font.italic = italic
    return txb


def arrow(slide, x1, y1, x2, y2, color=None, width=1.75, dashed=False):
    if color is None:
        color = QARROW
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    try:
        ln = conn.line._ln
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')
    except Exception:
        pass
    return conn


# ─── Build ───────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = blank(prs)
rect(s1, 0, 0, 13.33, 7.5, NAVY)
rect(s1, 0, 0, 0.22, 7.5, RGBColor(0x00, 0xB4, 0xF0))

lbl(s1, 0.45, 2.4,  12.0, 1.3,  'U1 Query Flow',
    fsize=52, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)
lbl(s1, 0.45, 3.85, 12.0, 0.9,
    '"How many UEs are currently active in the network?"',
    fsize=22, fc=RGBColor(0x99, 0xCC, 0xFF), italic=True, align=PP_ALIGN.LEFT)
lbl(s1, 0.45, 5.0,  9.0,  0.5,
    'Tool: query_network  →  sum connected_ues across all 30 cells',
    fsize=13, fc=RGBColor(0x55, 0xDD, 0x88), bold=True, align=PP_ALIGN.LEFT)
lbl(s1, 0.45, 5.7,  12.0, 0.4,
    '4G/5G NSA Network  ·  Malleswaram  ·  30 Cells  ·  3 DUs  ·  1 CU  ·  18,400 peak UEs',
    fsize=11, fc=RGBColor(0x77, 0x99, 0xBB), align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Flow Diagram
# ══════════════════════════════════════════════════════════════════════════════
s2 = blank(prs)
rect(s2, 0, 0, 13.33, 7.5, LGRAY)

# Header
rect(s2, 0, 0, 13.33, 0.9, NAVY)
lbl(s2, 0.3, 0.05, 9.5, 0.8, 'U1 Query Flow — System Data Path',
    fsize=22, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)
lbl(s2, 10.0, 0.1, 3.0, 0.7, 'query_network',
    fsize=15, fc=RGBColor(0x88, 0xCC, 0xFF), bold=True)

# Layer bands
rect(s2, 0.95, 1.0, 12.35, 1.8,  L_ORG,  border=RGBColor(0xDD, 0xAA, 0x80))
rect(s2, 0.95, 2.8, 12.35, 1.95, L_PUR,  border=RGBColor(0xAA, 0x88, 0xCC))
rect(s2, 0.95, 4.75,12.35, 2.65, L_TEAL, border=RGBColor(0x77, 0xBB, 0xBB))

# Layer labels (vertical strips on left)
lbl(s2, 0.0, 1.35, 0.98, 0.9,  'USER', fsize=8, fc=ORANGE, bold=True)
lbl(s2, 0.0, 3.1,  0.98, 0.9,  'AI',   fsize=8, fc=PURPLE, bold=True)
lbl(s2, 0.0, 5.4,  0.98, 0.9,  'DATA', fsize=8, fc=TEAL,   bold=True)

# ── Row 1: User Layer boxes ───────────────────────────────────────────────────
#   User         x=1.1  cx=1.85
box(s2, 1.1, 1.15, 1.5, 0.85, ORANGE, ['User', 'query input'], fsize=11)
#   chat.py      x=3.3  cx=4.25
box(s2, 3.3, 1.15, 1.9, 0.85, BLUE,   ['chat.py', 'localhost:8082'], fsize=11)

# ── Row 2: AI Layer boxes ─────────────────────────────────────────────────────
#   Orchestrator x=1.1  cx=2.2
box(s2, 1.1, 2.95, 2.2, 0.85, NAVY,   ['Orchestrator', 'port 8082'], fsize=11)
#   Gemini LLM   x=4.3  cx=5.65
box(s2, 4.3, 2.95, 2.6, 0.85, PURPLE, ['Gemini LLM', 'gemini-2.5-flash'], fsize=11)

# ── Row 3: Data Layer boxes ───────────────────────────────────────────────────
#   Controller   x=1.1  cx=2.2
box(s2, 1.1, 4.9,  2.2, 0.85, TEAL,  ['Controller', 'port 8080'], fsize=11)
#   topology.json x=4.1  cx=5.25
box(s2, 4.1, 5.85, 2.3, 0.85, GREEN, ['topology.json', 'cell/DU/CU config'], fsize=10)
#   InfluxDB      x=7.3  cx=8.45
box(s2, 7.3, 5.85, 2.2, 0.85, GREEN, ['InfluxDB', 'port 8086'], fsize=11)

# ── Arrows: Query path (BLUE) ─────────────────────────────────────────────────
# 1. User → chat.py
arrow(s2, 2.6, 1.575, 3.3, 1.575, color=QARROW)
lbl(s2, 2.6, 1.2, 0.7, 0.35, 'HTTP\nPOST', fsize=7, fc=BLUE)

# 2. chat.py → Orchestrator (diagonal: bottom of chat.py → top of Orchestrator)
arrow(s2, 4.25, 2.0,  2.2, 2.95, color=QARROW)
lbl(s2, 2.7, 2.35, 1.5, 0.35, 'POST /chat', fsize=7, fc=BLUE)

# 3. Orchestrator → Gemini (horizontal: right side → left side, top half)
arrow(s2, 3.3, 3.1, 4.3, 3.1, color=QARROW)
lbl(s2, 3.3, 2.75, 1.05, 0.35, 'prompt +\nsnapshot', fsize=7, fc=BLUE)

# 4. Gemini → Orchestrator (return, bottom half, GREEN)
arrow(s2, 4.3, 3.55, 3.3, 3.55, color=RARROW)
lbl(s2, 3.3, 3.6, 1.05, 0.3, 'fn call:\nquery_network', fsize=7, fc=GREEN)

# 5. Orchestrator → Controller (vertical down)
arrow(s2, 2.2, 3.8, 2.2, 4.9, color=QARROW)
lbl(s2, 2.25, 4.3, 1.2, 0.35, 'GET /network', fsize=7, fc=BLUE)

# 6. Controller → topology.json (diagonal)
arrow(s2, 2.5, 5.75, 4.1, 6.1, color=QARROW)
lbl(s2, 2.7, 5.55, 1.3, 0.3, 'read config', fsize=7, fc=BLUE)

# 7. Controller → InfluxDB (diagonal)
arrow(s2, 3.0, 5.6, 7.3, 5.95, color=QARROW)
lbl(s2, 4.7, 5.35, 2.0, 0.3, 'Flux query: connected_ues', fsize=7, fc=BLUE)

# ── Arrows: Response path (GREEN, dashed) ────────────────────────────────────
# 8. topology.json → Controller
arrow(s2, 4.6, 5.85, 3.1, 5.0, color=RARROW, dashed=True)
lbl(s2, 3.6, 5.3, 1.1, 0.3, 'cell config', fsize=7, fc=GREEN)

# 9. InfluxDB → Controller
arrow(s2, 7.5, 5.85, 3.2, 5.2, color=RARROW, dashed=True)
lbl(s2, 5.0, 5.65, 2.1, 0.3, 'connected_ues per cell', fsize=7, fc=GREEN)

# 10. Controller → Orchestrator (vertical up, return)
arrow(s2, 1.9, 4.9, 1.9, 3.8, color=RARROW, dashed=True)
lbl(s2, 0.1, 4.2, 1.75, 0.45, 'merged\nsnapshot', fsize=7, fc=GREEN)

# 11. Orchestrator → chat.py → User (return, horizontal)
arrow(s2, 3.3, 3.65, 2.3, 2.1, color=RARROW, dashed=True)
arrow(s2, 2.25, 2.05, 1.85, 1.75, color=RARROW, dashed=True)
lbl(s2, 0.6, 2.5, 1.4, 0.35, 'NL reply', fsize=7, fc=GREEN)

# ── Legend ────────────────────────────────────────────────────────────────────
rect(s2, 10.0, 1.15, 3.1, 1.45, WHITE, border=RGBColor(0xCC, 0xCC, 0xCC))
lbl(s2, 10.05, 1.18, 3.0, 0.35, 'Legend', fsize=10, fc=NAVY, bold=True)
arrow(s2, 10.1, 1.72, 10.7, 1.72, color=QARROW, width=1.5)
lbl(s2, 10.75, 1.58, 2.2, 0.3, 'Query / Request path', fsize=9, fc=DGRAY)
arrow(s2, 10.1, 2.09, 10.7, 2.09, color=RARROW, width=1.5, dashed=True)
lbl(s2, 10.75, 1.96, 2.2, 0.3, 'Response / Return path', fsize=9, fc=DGRAY)

# Gemini note box
box(s2, 7.8, 2.95, 3.0, 0.85, RGBColor(0x44, 0x22, 0x70),
    ['Sums connected_ues', 'across all 30 cells', '→ formulates reply'],
    fsize=9, bold0=False)
arrow(s2, 6.9, 3.375, 7.8, 3.375, color=PURPLE, width=1.5)
lbl(s2, 6.9, 3.1, 0.9, 0.3, 'tool\nresult', fsize=7, fc=PURPLE)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Step-by-Step Table
# ══════════════════════════════════════════════════════════════════════════════
s3 = blank(prs)
rect(s3, 0, 0, 13.33, 7.5, LGRAY)
rect(s3, 0, 0, 13.33, 0.9, NAVY)
lbl(s3, 0.3, 0.05, 10.0, 0.8, 'U1 Query — Step-by-Step Trace',
    fsize=22, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)

steps = [
    ('1', 'User',          'Types query into chat.py'),
    ('2', 'chat.py',       'HTTP POST {session, message} to Orchestrator /chat'),
    ('3', 'Orchestrator',  'Injects live network snapshot into system prompt; calls Gemini'),
    ('4', 'Gemini LLM',    'Decides query_network is the correct tool; emits function call'),
    ('5', 'Orchestrator',  'Dispatches query_network → GET /network on Controller'),
    ('6', 'Controller',    'Reads topology.json for cell/DU/CU configuration'),
    ('7', 'Controller',    'Runs Flux query against InfluxDB: latest connected_ues per cell'),
    ('8', 'Controller',    'Merges topology + KPIs; returns full 30-cell network snapshot'),
    ('9', 'Orchestrator',  'Passes tool result back to Gemini'),
    ('10','Gemini LLM',    'Sums connected_ues across all cells to get network-wide total'),
    ('11','Gemini LLM',    'Generates natural language answer'),
    ('12','Orchestrator',  'Returns HTTP response {reply} to chat.py'),
    ('13','chat.py',       'Prints answer to user'),
]

col_colors = {
    'User':         ORANGE,
    'chat.py':      BLUE,
    'Orchestrator': NAVY,
    'Gemini LLM':   PURPLE,
    'Controller':   TEAL,
}

ROW_H = 0.46
START_Y = 1.0
for i, (step, actor, action) in enumerate(steps):
    y = START_Y + i * ROW_H
    bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xE8, 0xEC, 0xF2)
    rect(s3, 0.3, y, 12.7, ROW_H, bg, border=RGBColor(0xCC, 0xCC, 0xCC), bw=0.3)

    # Step number
    lbl(s3, 0.35, y + 0.05, 0.55, ROW_H - 0.1, step,
        fsize=11, fc=NAVY, bold=True)

    # Actor badge
    actor_color = col_colors.get(actor, DGRAY)
    rect(s3, 0.95, y + 0.08, 1.6, ROW_H - 0.16, actor_color,
         border=RGBColor(0xAA, 0xAA, 0xAA), bw=0.3)
    lbl(s3, 0.95, y + 0.08, 1.6, ROW_H - 0.16, actor,
        fsize=9, fc=WHITE, bold=True)

    # Action text
    lbl(s3, 2.65, y + 0.05, 10.3, ROW_H - 0.1, action,
        fsize=10, fc=DGRAY, align=PP_ALIGN.LEFT)


# ─── Save ─────────────────────────────────────────────────────────────────────
out = r'c:\Users\gurs2\OneDrive\Documents\GitHub\telecom-automation\test\U1_Query_Flow.pptx'
prs.save(out)
print(f'Saved: {out}')
