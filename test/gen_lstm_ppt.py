"""
Generate LSTM_Explained.pptx — 10-slide visual explainer for the KPI LSTM classifier.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LSTM_Explained.pptx")

# ── colours ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x29, 0x55)
BLUE   = RGBColor(0x17, 0x5C, 0xA6)
TEAL   = RGBColor(0x00, 0x7E, 0x8A)
GREEN  = RGBColor(0x1B, 0x6F, 0x42)
ORANGE = RGBColor(0xC4, 0x50, 0x08)
PURPLE = RGBColor(0x5A, 0x23, 0x8C)
RED    = RGBColor(0xC0, 0x00, 0x00)
AMBER  = RGBColor(0xFF, 0xC0, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
DGRAY  = RGBColor(0x59, 0x59, 0x59)

C_NORMAL  = RGBColor(0x1B, 0x6F, 0x42)
C_OVER    = RGBColor(0xC4, 0x50, 0x08)
C_UNDER   = RGBColor(0x17, 0x5C, 0xA6)
C_SINR    = RGBColor(0xC0, 0x00, 0x00)
C_PWR     = RGBColor(0x5A, 0x23, 0x8C)

BG_NORMAL = RGBColor(0xC6, 0xEF, 0xCE)
BG_OVER   = RGBColor(0xFF, 0xEB, 0x9C)
BG_UNDER  = RGBColor(0xBD, 0xD7, 0xEE)
BG_SINR   = RGBColor(0xFF, 0xC7, 0xCE)
BG_PWR    = RGBColor(0xE2, 0xEF, 0xDA)


# ── helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def rgb(r, g, b):
    return RGBColor(r, g, b)


def rect(slide, x, y, w, h, fill, text="", fs=10, bold=False,
         fc=WHITE, align=PP_ALIGN.CENTER, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.color.rgb = fc
    return shape


def box(slide, x, y, w, h, fill, text="", fs=10, bold=False,
        fc=WHITE, align=PP_ALIGN.CENTER, line_color=None):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    lc = line_color or fill
    shape.line.color.rgb = lc
    shape.line.width = Pt(1.2)
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(fs)
        run.font.bold = bold
        run.font.color.rgb = fc
    return shape


def lbl(slide, x, y, w, h, text, fs=9, fc=NAVY, bold=False,
        align=PP_ALIGN.CENTER, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fs)
    run.font.bold = bold
    run.font.color.rgb = fc
    run.font.italic = italic
    return tb


def arr(slide, x1, y1, x2, y2, color=NAVY, width=Pt(2)):
    from pptx.util import Emu
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    cx = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cx.line.color.rgb = color
    cx.line.width = width
    ln = cx.line._get_or_add_ln()
    etree.SubElement(ln, qn('a:tailEnd'), attrib={'type': 'arrow'})
    return cx


def title_bar(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 0.72, NAVY, title, fs=22, bold=True, fc=WHITE)
    if subtitle:
        rect(slide, 0, 0.72, 13.33, 0.32, BLUE, subtitle, fs=10, fc=WHITE)


def slide_num(slide, n):
    lbl(slide, 12.8, 7.1, 0.4, 0.3, str(n), fs=8, fc=DGRAY)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
def slide_title(prs):
    s = blank(prs)
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    # accent bar
    rect(s, 0, 3.2, 13.33, 0.06, TEAL)

    lbl(s, 1.0, 1.2, 11.33, 1.0,
        "LSTM KPI Classifier", fs=42, fc=WHITE, bold=True)
    lbl(s, 1.0, 2.3, 11.33, 0.7,
        "Explained Simply", fs=28, fc=RGBColor(0xA0, 0xC8, 0xFF))
    lbl(s, 1.0, 3.5, 11.33, 0.5,
        "How AI watches 30 cells every 10 seconds and keeps the Malleswaram network healthy",
        fs=13, fc=WHITE, italic=True)

    # 5 class pills at bottom
    classes = [
        ("NORMAL",      BG_NORMAL, C_NORMAL),
        ("OVERLOAD",    BG_OVER,   C_OVER),
        ("UNDERLOAD",   BG_UNDER,  C_UNDER),
        ("SINR LOW",    BG_SINR,   C_SINR),
        ("POWER WASTE", BG_PWR,    C_PWR),
    ]
    for i, (name, bg, fg) in enumerate(classes):
        box(s, 0.5 + i * 2.55, 6.5, 2.3, 0.55, bg, name, fs=11, bold=True, fc=fg)

    lbl(s, 0.5, 7.15, 12.0, 0.3,
        "agents/kpi_agent/model.py  |  train.py  |  kpi_agent.py",
        fs=8, fc=RGBColor(0x80, 0xA0, 0xC0), italic=True)
    slide_num(s, 1)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — The Problem: Why Not Just Rules?
# ─────────────────────────────────────────────────────────────────────────────
def slide_problem(prs):
    s = blank(prs)
    title_bar(s, "The Problem — Why Not Just a Threshold Rule?",
              "A single reading misses trends. The LSTM reads 60 seconds of history.")
    slide_num(s, 2)

    # Left: Rule approach
    rect(s, 0.3, 1.2, 5.8, 0.55, RED, "RULE-BASED  (current fallback)", fs=12, bold=True)
    box(s, 0.3, 1.85, 5.8, 3.8, RGBColor(0xFF, 0xF0, 0xF0),
        fc=NAVY, line_color=RED)

    lbl(s, 0.5, 2.0, 5.4, 0.4,
        "IF  prb_dl_pct > 85%  THEN  OVERLOAD", fs=11, bold=True, fc=RED)
    lbl(s, 0.5, 2.5, 5.4, 2.8,
        "Looks at ONE reading at t=now\n\n"
        "PRB = 86%  at  t=now\n"
        "  → fires OVERLOAD alert immediately\n\n"
        "But what if PRB was 87% last second\n"
        "and is now DROPPING back to 70%?\n\n"
        "   False alarm — no real problem!",
        fs=10, fc=DGRAY, align=PP_ALIGN.LEFT)

    box(s, 0.3, 5.75, 5.8, 0.7,
        RGBColor(0xFF, 0xC7, 0xCE),
        "Misses trends  ·  False alarms  ·  Acts too late",
        fs=10, bold=True, fc=RED, line_color=RED)

    # Right: LSTM approach
    rect(s, 7.2, 1.2, 5.8, 0.55, GREEN, "LSTM  (AI model)", fs=12, bold=True)
    box(s, 7.2, 1.85, 5.8, 3.8, RGBColor(0xF0, 0xFF, 0xF4),
        fc=NAVY, line_color=GREEN)

    lbl(s, 7.4, 2.0, 5.4, 0.4,
        "Reads the last 6 readings (60 seconds)", fs=11, bold=True, fc=GREEN)

    # PRB trend table
    for i, (t, prb, clr) in enumerate([
        ("t-50s", "60%", LGRAY), ("t-40s", "68%", LGRAY),
        ("t-30s", "75%", BG_OVER), ("t-20s", "82%", BG_OVER),
        ("t-10s", "89%", BG_SINR), ("t=now", "94%", BG_SINR),
    ]):
        rect(s, 7.4 + i * 0.88, 2.55, 0.82, 0.38, clr,
             t, fs=7, fc=NAVY, bold=(i >= 4))
        rect(s, 7.4 + i * 0.88, 2.98, 0.82, 0.38,
             clr if i < 4 else BG_SINR,
             prb, fs=10, bold=(i >= 4),
             fc=RED if i >= 4 else NAVY)

    lbl(s, 7.4, 3.5, 5.4, 2.2,
        "Trend: PRB rising consistently for 60s\n\n"
        "  Predicts OVERLOAD BEFORE it peaks\n\n"
        "  Acts while there is still time to\n"
        "  move the cell to a lighter DU",
        fs=10, fc=DGRAY, align=PP_ALIGN.LEFT)

    box(s, 7.2, 5.75, 5.8, 0.7,
        BG_NORMAL,
        "Catches trends early  ·  Fewer false alarms  ·  Proactive action",
        fs=10, bold=True, fc=GREEN, line_color=GREEN)

    # VS label
    lbl(s, 6.2, 3.2, 0.9, 0.6, "VS", fs=22, bold=True, fc=NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — What the LSTM Sees: 9 Features × 6 Timesteps
# ─────────────────────────────────────────────────────────────────────────────
def slide_input(prs):
    s = blank(prs)
    title_bar(s, "What the LSTM Sees — 9 KPI Features × 6 Timesteps (60 seconds)",
              "Each cell gets one reading every 10 s  →  6 readings = 1 input tensor  (shape: 1 × 6 × 9)")
    slide_num(s, 3)

    features = [
        ("PRB DL %",       "Physical Resource Block usage",        "0 – 100 %"),
        ("SINR dB",        "Signal quality (higher = better)",      "-5 – +30 dB"),
        ("Connected UEs",  "Active users on this cell",            "0 – 800"),
        ("Power W",        "RF transmit power draw",               "0 – 1200 W"),
        ("Packet Loss %",  "Percentage of lost data packets",      "0 – 5 %"),
        ("Throughput Mbps","Downlink data speed delivered",        "0 – 4000 Mbps"),
        ("CQI",            "Channel Quality Indicator (link quality)","0 – 15"),
        ("BLER %",         "Block Error Rate (retransmissions)",   "0 – 30 %"),
        ("Latency ms",     "Round-trip cell latency",              "0 – 500 ms"),
    ]

    times = ["t-50s", "t-40s", "t-30s", "t-20s", "t-10s", "t=now"]
    col_w = 1.55
    row_h = 0.52
    x0, y0 = 0.25, 1.18

    # Header row (time labels)
    rect(s, x0, y0, 3.5, row_h, NAVY, "Feature / Timestep",
         fs=9, bold=True, fc=WHITE)
    for ci, t in enumerate(times):
        clr = TEAL if t == "t=now" else BLUE
        rect(s, x0 + 3.5 + ci * col_w, y0, col_w, row_h,
             clr, t, fs=9, bold=True, fc=WHITE)

    # Sample values for a NORMAL-trending cell
    sample = [
        ["52%",  "55%",  "58%",  "57%",  "54%",  "53%"],
        ["17.2", "16.8", "17.1", "17.5", "16.9", "17.3"],
        ["312",  "318",  "325",  "320",  "315",  "322"],
        ["491",  "495",  "502",  "498",  "493",  "496"],
        ["0.03", "0.04", "0.03", "0.05", "0.04", "0.03"],
        ["1230", "1260", "1285", "1270", "1240", "1255"],
        ["11",   "11",   "10",   "11",   "11",   "11"],
        ["1.4",  "1.5",  "1.3",  "1.5",  "1.4",  "1.4"],
        ["12",   "13",   "12",   "11",   "13",   "12"],
    ]

    for ri, (feat, desc, rng) in enumerate(features):
        yy = y0 + (ri + 1) * row_h
        bg = LGRAY if ri % 2 == 0 else WHITE
        rect(s, x0, yy, 3.5, row_h, bg,
             f"{feat}\n({rng})", fs=8, fc=NAVY, bold=(ri == 0),
             align=PP_ALIGN.LEFT)
        for ci, val in enumerate(sample[ri]):
            cell_bg = TEAL if ci == 5 else bg
            cell_fc = WHITE if ci == 5 else NAVY
            rect(s, x0 + 3.5 + ci * col_w, yy, col_w, row_h,
                 cell_bg, val, fs=9, bold=(ci == 5), fc=cell_fc)

    # Annotation
    lbl(s, 0.25, 6.85, 9.0, 0.4,
        "All 54 values (9 features × 6 timesteps) are normalised to [0, 1] before feeding to the LSTM",
        fs=9, fc=DGRAY, italic=True)

    box(s, 9.6, 1.18, 3.5, 5.9, RGBColor(0xF0, 0xF4, 0xFF),
        fc=NAVY, line_color=BLUE)
    lbl(s, 9.7, 1.3, 3.3, 0.4, "model.py constants", fs=9, bold=True, fc=BLUE)
    lbl(s, 9.7, 1.75, 3.3, 5.2,
        "SEQ_LEN  = 6\n(6 timesteps)\n\n"
        "N_FEATURES = 9\n(9 KPI fields)\n\n"
        "N_CLASSES = 5\n(output labels)\n\n"
        "Input shape:\n(1, 6, 9)\nbatch × time × feat",
        fs=10, fc=NAVY, align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — Normalisation
# ─────────────────────────────────────────────────────────────────────────────
def slide_norm(prs):
    s = blank(prs)
    title_bar(s, "Normalisation — Making All Values Comparable",
              "Raw KPI values have very different scales. Normalise everything to 0–1 before the LSTM sees it.")
    slide_num(s, 4)

    # Formula box
    box(s, 0.3, 1.15, 12.7, 0.75, NAVY,
        "normalised  =  ( value  -  minimum )  /  range",
        fs=16, bold=True, fc=WHITE)

    # Examples table
    headers = ["Feature", "Raw Value", "Min", "Range", "Normalised", "Meaning"]
    widths  = [2.5, 1.7, 1.2, 1.2, 1.8, 4.2]
    x0, y0 = 0.3, 2.1

    xpos = x0
    for h, w in zip(headers, widths):
        rect(s, xpos, y0, w, 0.42, BLUE, h, fs=10, bold=True, fc=WHITE)
        xpos += w

    rows = [
        ("PRB DL %",       "94",    "0",   "100",  "0.94",  "Near full capacity"),
        ("SINR dB",        "2",     "-5",  "35",   "0.20",  "Poor signal quality"),
        ("Connected UEs",  "720",   "0",   "800",  "0.90",  "Heavily loaded"),
        ("Power W",        "880",   "0",   "1200", "0.73",  "High but not peak"),
        ("Pkt Loss %",     "0.85",  "0",   "5",    "0.17",  "Some congestion"),
        ("Throughput Mbps","3050",  "0",   "4000", "0.76",  "Good throughput"),
        ("CQI",            "7",     "0",   "15",   "0.47",  "Degraded channel"),
        ("BLER %",         "8.0",   "0",   "30",   "0.27",  "Elevated errors"),
        ("Latency ms",     "38",    "0",   "500",  "0.08",  "Normal latency"),
    ]
    for ri, row in enumerate(rows):
        yy = y0 + (ri + 1) * 0.43
        bg = BG_OVER if ri <= 1 else (LGRAY if ri % 2 == 0 else WHITE)
        xpos = x0
        for val, w in zip(row, widths):
            fc = RED if ri <= 1 and w <= 1.8 else NAVY
            bold = ri <= 1
            rect(s, xpos, yy, w, 0.43, bg, val, fs=9,
                 bold=bold, fc=fc)
            xpos += w

    lbl(s, 0.3, 6.75, 13.0, 0.4,
        "Why? — PRB is 0-100, Power is 0-1200, Latency is 0-500. "
        "Without normalisation the LSTM would learn 'power matters more' just because its numbers are bigger.",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — Inside the LSTM: Gates and Memory
# ─────────────────────────────────────────────────────────────────────────────
def slide_gates(prs):
    s = blank(prs)
    title_bar(s, "Inside One LSTM Cell — Gates and Memory",
              "Three gates decide what to remember, what to forget, and what to pass forward at each timestep")
    slide_num(s, 5)

    # Cell state bar (long-term memory) at top
    rect(s, 0.5, 1.1, 12.3, 0.6, TEAL,
         "CELL STATE  —  Long-Term Memory   "
         "(carries information across all 6 timesteps without vanishing)",
         fs=11, bold=True, fc=WHITE)

    # Three gates
    gates = [
        ("FORGET GATE", ORANGE,
         "What from the past\nis no longer relevant?",
         "Example:\n\"Last hour's off-peak\nload data — forget it.\nCurrent peak hour\nis what matters.\""),
        ("INPUT GATE", TEAL,
         "What new information\nshould I remember?",
         "Example:\n\"PRB just jumped\nfrom 60% to 89%.\nThat spike is important\n— store it.\""),
        ("OUTPUT GATE", PURPLE,
         "What do I pass\nto the next timestep?",
         "Example:\n\"Rising PRB trend\nover 5 timesteps\n— send that signal\nforward.\""),
    ]
    for i, (name, clr, q, ex) in enumerate(gates):
        gx = 0.5 + i * 4.15
        rect(s, gx, 1.85, 3.8, 0.48, clr, name, fs=12, bold=True, fc=WHITE)
        box(s, gx, 2.38, 3.8, 1.3, LGRAY, q, fs=10, fc=NAVY, line_color=clr)
        box(s, gx, 3.75, 3.8, 1.9, WHITE, ex, fs=9, fc=DGRAY,
            align=PP_ALIGN.LEFT, line_color=clr)
        arr(s, gx + 1.9, 3.68, gx + 1.9, 3.78, color=clr)

    # Hidden state bar at bottom
    rect(s, 0.5, 5.75, 12.3, 0.55, BLUE,
         "HIDDEN STATE  —  Short-Term Memory   "
         "(output of this timestep, passed as input to the next)",
         fs=11, bold=True, fc=WHITE)

    # Arrows from gates to cell/hidden state
    for x in [2.38, 6.53, 10.68]:
        arr(s, x, 5.70, x, 5.78, color=NAVY)

    lbl(s, 0.5, 6.4, 12.3, 0.4,
        "The CELL STATE solves the 'vanishing gradient' problem — it carries signals across many timesteps without losing them, "
        "so the LSTM can learn from events that happened 50 seconds ago.",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — Bidirectional LSTM
# ─────────────────────────────────────────────────────────────────────────────
def slide_bidir(prs):
    s = blank(prs)
    title_bar(s, "Bidirectional — The LSTM Reads Forward AND Backward",
              "Two LSTMs run in parallel. Their outputs are concatenated → richer context.")
    slide_num(s, 6)

    # Timestep boxes
    ts = ["t-50s\nPRB 60%", "t-40s\nPRB 68%", "t-30s\nPRB 75%",
          "t-20s\nPRB 82%", "t-10s\nPRB 89%", "t=now\nPRB 94%"]
    colors_ts = [LGRAY, LGRAY, BG_OVER, BG_OVER, BG_SINR, BG_SINR]
    for i, (t, c) in enumerate(zip(ts, colors_ts)):
        fc = RED if i >= 4 else NAVY
        rect(s, 0.4 + i * 2.08, 1.1, 1.9, 0.7, c, t,
             fs=9, bold=(i >= 4), fc=fc)

    # Forward LSTM row
    rect(s, 0.4, 2.05, 12.5, 0.55, TEAL,
         "Forward LSTM  →  reads t-50s ... t=now", fs=11, bold=True, fc=WHITE)
    lbl(s, 0.4, 2.65, 12.5, 0.4,
        '"PRB is climbing steadily — this looks like sustained overload building up"',
        fs=10, fc=TEAL, italic=True)

    arr(s, 0.8, 1.82, 12.6, 1.82, color=TEAL, width=Pt(2.5))

    # Backward LSTM row
    rect(s, 0.4, 3.25, 12.5, 0.55, PURPLE,
         "Backward LSTM  ←  reads t=now ... t-50s", fs=11, bold=True, fc=WHITE)
    lbl(s, 0.4, 3.85, 12.5, 0.4,
        '"At t=now PRB=94% (confirmed peak). Earlier values confirm the trend was real, not a spike"',
        fs=10, fc=PURPLE, italic=True)

    arr(s, 12.6, 3.22, 0.8, 3.22, color=PURPLE, width=Pt(2.5))

    # Combine
    rect(s, 2.0, 4.5, 9.3, 0.58, NAVY,
         "Concatenate forward + backward outputs  →  128-dimensional vector  "
         "(64 forward + 64 backward hidden units)",
         fs=11, bold=True, fc=WHITE)

    box(s, 3.5, 5.25, 6.3, 1.1, BG_NORMAL, line_color=GREEN,
        fc=NAVY)
    lbl(s, 3.5, 5.3, 6.3, 1.0,
        "Result:  OVERLOAD predicted at 97% confidence\n"
        "The backward pass confirmed the trend was real, not noise",
        fs=11, bold=True, fc=GREEN)

    arr(s, 6.67, 5.08, 6.67, 5.27, color=NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 7 — Full Architecture
# ─────────────────────────────────────────────────────────────────────────────
def slide_arch(prs):
    s = blank(prs)
    title_bar(s, "Full Architecture — From Raw KPIs to a Classification Decision",
              "Two stacked bidirectional LSTM layers followed by a two-layer MLP head")
    slide_num(s, 7)

    layers = [
        (BLUE,   "INPUT TENSOR",
         "Shape: (1, 6, 9)  —  1 cell  ×  6 timesteps  ×  9 normalised KPI features"),
        (TEAL,   "BiLSTM LAYER 1   (hidden=64, bidirectional)",
         "Reads each of the 6 timesteps in both directions  →  output shape (1, 6, 128)   [ 64 fwd + 64 bwd ]"),
        (TEAL,   "BiLSTM LAYER 2   (hidden=64, bidirectional, dropout 25%)",
         "Deeper temporal abstraction  →  output shape (1, 6, 128)  ·  dropout prevents overfitting"),
        (PURPLE, "TAKE LAST TIMESTEP  out[:, -1, :]",
         "Extract only the final timestep's combined state  →  shape (1, 128)  ·  this is the 'summary' of the whole 60s"),
        (ORANGE, "LINEAR 128 → 64  +  ReLU  +  Dropout 25%",
         "First MLP layer: compress and non-linearise the LSTM output"),
        (ORANGE, "LINEAR 64 → 5",
         "Second MLP layer: project to 5 output scores (one per class)"),
        (GREEN,  "SOFTMAX  →  Probabilities",
         "Convert raw scores to probabilities that sum to 1.0   e.g. [NORMAL 0.04, OVERLOAD 0.93, ...]"),
    ]

    for i, (clr, title, desc) in enumerate(layers):
        y = 1.1 + i * 0.84
        rect(s, 0.3, y, 3.5, 0.72, clr, title, fs=9, bold=True, fc=WHITE)
        rect(s, 3.85, y, 9.1, 0.72, LGRAY if i % 2 == 0 else WHITE,
             desc, fs=9, fc=NAVY, align=PP_ALIGN.LEFT)
        if i < len(layers) - 1:
            arr(s, 1.9, y + 0.72, 1.9, y + 0.84, color=clr)

    # code snippet on right
    box(s, 3.85, 1.1, 9.1, 0.0, WHITE, line_color=WHITE)  # placeholder

    lbl(s, 0.3, 7.0, 12.7, 0.38,
        "model.py  KPIClassifier  ·  Total parameters ≈ 145 K  ·  "
        "Trains in ~2 minutes on CPU (50 epochs, 5 000 synthetic samples)",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 8 — 5 Output Classes
# ─────────────────────────────────────────────────────────────────────────────
def slide_classes(prs):
    s = blank(prs)
    title_bar(s, "5 Output Classes — What Decision Does the LSTM Make?",
              "Softmax outputs a probability for each class. The highest probability wins. MIN_CONFIDENCE = 70%.")
    slide_num(s, 8)

    classes = [
        (0, "NORMAL",      C_NORMAL, BG_NORMAL, "70% of real cells",
         "PRB 40-70%  ·  Good SINR  ·  Stable UEs",
         "No action needed.\nNetwork is healthy."),
        (1, "OVERLOAD",    C_OVER,   BG_OVER,   "15% of real cells",
         "PRB >85%  ·  Packet loss rising  ·  Many UEs",
         "Auto-move cell to\nlightest available DU"),
        (2, "UNDERLOAD",   C_UNDER,  BG_UNDER,  "8% of real cells",
         "PRB <20%  ·  Very few UEs  ·  Wasted capacity",
         "Flag as sleep candidate\n(energy saving / DTX)"),
        (3, "SINR LOW",    C_SINR,   BG_SINR,   "5% of real cells",
         "SINR <5 dB  ·  High BLER  ·  Low CQI",
         "CRITICAL alert +\nrequest PCI re-optimisation"),
        (4, "POWER WASTE", C_PWR,    BG_PWR,    "2% of real cells",
         "High power draw  ·  Very few UEs (<15)",
         "WARNING alert +\nrecommend DTX / sleep mode"),
    ]

    for i, (idx, name, fg, bg, freq, kpis, action) in enumerate(classes):
        y = 1.1 + i * 1.22
        # Index badge
        rect(s, 0.3, y, 0.55, 0.95, fg,
             str(idx), fs=20, bold=True, fc=WHITE)
        # Class name
        rect(s, 0.9, y, 2.8, 0.95, bg, name, fs=14, bold=True, fc=fg)
        # Frequency
        rect(s, 3.75, y, 2.0, 0.95, LGRAY, freq, fs=10, fc=DGRAY)
        # KPI signature
        rect(s, 5.8, y, 4.1, 0.95, WHITE,
             kpis, fs=9, fc=NAVY, align=PP_ALIGN.LEFT)
        # Action
        rect(s, 9.95, y, 3.1, 0.95, bg,
             action, fs=10, bold=True, fc=fg)

    # Column headers
    for x, w, h in [(0.3, 0.55, ""), (0.9, 2.8, "Class"),
                    (3.75, 2.0, "Frequency"), (5.8, 4.1, "Typical KPI Signature"),
                    (9.95, 3.1, "SON Action")]:
        if h:
            rect(s, x, 0.85, w, 0.25, BLUE, h, fs=8, bold=True, fc=WHITE)

    lbl(s, 0.3, 7.12, 12.7, 0.3,
        "Frequency is the realistic class distribution used in training (train.py CLASS_COUNTS). "
        "Weighted sampler ensures rare classes are learned even at 2%.",
        fs=8, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 9 — Training Data
# ─────────────────────────────────────────────────────────────────────────────
def slide_training(prs):
    s = blank(prs)
    title_bar(s, "Training — How the Model Was Taught",
              "Synthetic data generated from two separate profiles: 5G NR (64T64R) and 4G LTE (4T4R)  ·  60 epochs")
    slide_num(s, 9)

    # Class distribution bars
    lbl(s, 0.3, 1.1, 4.0, 0.38, "Class Distribution", fs=12, bold=True, fc=NAVY)
    classes_dist = [
        ("NORMAL",      3500, BG_NORMAL, C_NORMAL),
        ("OVERLOAD",     750, BG_OVER,   C_OVER),
        ("UNDERLOAD",    400, BG_UNDER,  C_UNDER),
        ("SINR LOW",     250, BG_SINR,   C_SINR),
        ("POWER WASTE",  100, BG_PWR,    C_PWR),
    ]
    total = 5000
    for i, (name, n, bg, fg) in enumerate(classes_dist):
        y = 1.55 + i * 0.72
        rect(s, 0.3, y, 1.8, 0.55, bg, name, fs=10, bold=True, fc=fg)
        bw = (n / total) * 7.0
        rect(s, 2.15, y, bw, 0.55, fg, f"{n:,}  ({n*100//total}%)",
             fs=10, bold=True, fc=WHITE)

    # 5G vs 4G profiles panel
    lbl(s, 9.5, 1.1, 3.6, 0.38, "Two Sub-profiles per Class", fs=12, bold=True, fc=NAVY)

    rect(s, 9.5, 1.5, 1.7, 0.42, TEAL, "5G n78", fs=10, bold=True, fc=WHITE)
    rect(s, 11.25, 1.5, 1.8, 0.42, BLUE, "4G B3/B40", fs=10, bold=True, fc=WHITE)

    rows_5g4g = [
        ("NORMAL power", "520 W", "120 W"),
        ("OVERLOAD UEs", "720",   "230"),
        ("UNDERLOAD PRB", "9%",   "8%"),
        ("SINR_LOW BLER", "12%",  "10%"),
        ("PW_WASTE power", "880 W", "175 W"),
    ]
    for i, (feat, v5g, v4g) in enumerate(rows_5g4g):
        yy = 1.98 + i * 0.52
        bg = LGRAY if i % 2 == 0 else WHITE
        rect(s, 9.5, yy, 1.7, 0.5, bg, feat, fs=8, fc=NAVY, align=PP_ALIGN.LEFT)
        rect(s, 11.25, yy, 0.88, 0.5, bg, v5g, fs=9, fc=TEAL, bold=True)
        rect(s, 12.18, yy, 0.88, 0.5, bg, v4g, fs=9, fc=BLUE, bold=True)

    # Training process
    steps_tr = [
        (BLUE,   "Generate",  "5 000 sequences\n(5G + 4G)\n6 steps each"),
        (TEAL,   "Shuffle",   "Random permutation\nto mix classes\nand technologies"),
        (PURPLE, "Sample",    "WeightedRandom-\nSampler balances\nrare classes"),
        (ORANGE, "Train",     "60 epochs\nAdam + Cosine\nlr annealing"),
        (GREEN,  "Save",      "kpi_model.pt\nLoaded at\ncontainer start"),
    ]
    for i, (clr, title, body) in enumerate(steps_tr):
        x = 0.3 + i * 1.85
        rect(s, x, 5.5, 1.7, 0.42, clr, title, fs=10, bold=True, fc=WHITE)
        box(s, x, 5.97, 1.7, 1.2, LGRAY, body, fs=9, fc=NAVY, line_color=clr)
        if i < len(steps_tr) - 1:
            arr(s, x + 1.7, 5.72, x + 1.85, 5.72, color=clr)

    lbl(s, 0.3, 7.22, 12.7, 0.3,
        "Temporal drift added between timesteps so the LSTM learns trends, not just point-in-time snapshots  "
        "|  80/20 train/val split  |  Gradient clipping at 1.0",
        fs=8, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 10 — Live Pipeline: Cell to Action
# ─────────────────────────────────────────────────────────────────────────────
def slide_pipeline(prs):
    s = blank(prs)
    title_bar(s, "Live Pipeline — From a Cell Reading to an Automated SON Action",
              "Every 10 seconds per cell  ·  30 cells  ·  first 60 s uses rule-based fallback")
    slide_num(s, 10)

    steps = [
        (BLUE,   "DU Simulator\n(or real DU)",
         "Pushes cell_kpi\nto InfluxDB\nevery 10 s"),
        (TEAL,   "InfluxDB",
         "Stores time-series\nKPI telemetry\nfor all 30 cells"),
        (ORANGE, "KPI Agent\npolls /query",
         "Reads latest\n9 KPI fields\nfor each cell"),
        (PURPLE, "Sliding Buffer\n(deque maxlen=6)",
         "Appends reading.\nFull after 60 s.\nOlder readings drop off."),
        (TEAL,   "Normalise\n+ LSTM Infer",
         "9 raw values\n→ normalised\n→ model(x)"),
        (GREEN,  "Softmax\nProbabilities",
         "[0.04, 0.93,\n 0.01, 0.01,\n 0.01]\nconf=93%"),
        (ORANGE, "SON Action",
         "OVERLOAD:\nMove cell to\nlightest DU"),
    ]

    bw = 1.72
    for i, (clr, title, body) in enumerate(steps):
        x = 0.25 + i * bw
        rect(s, x, 1.1, 1.6, 0.52, clr, title, fs=9, bold=True, fc=WHITE)
        box(s, x, 1.67, 1.6, 1.1, LGRAY, body, fs=9, fc=NAVY, line_color=clr)
        if i < len(steps) - 1:
            arr(s, x + 1.6, 1.36, x + bw, 1.36, color=clr, width=Pt(2))

    # Confidence gate
    box(s, 3.5, 3.1, 6.3, 0.65, BG_OVER, line_color=ORANGE)
    lbl(s, 3.5, 3.12, 6.3, 0.6,
        "Confidence gate: if conf < 70% (MIN_CONFIDENCE) — log only, do NOT act",
        fs=10, bold=True, fc=ORANGE)

    # Fallback path
    rect(s, 0.25, 3.1, 3.0, 0.65, BLUE,
         "First 60 s: Rule-based fallback\n(buffer not full yet)", fs=9, bold=True, fc=WHITE)
    arr(s, 1.75, 2.78, 1.75, 3.12, color=BLUE)

    # SON actions detail
    son_actions = [
        (C_OVER,  BG_OVER,  "OVERLOAD",    "Move cell to lightest DU\nCooldown: 30 s between moves"),
        (C_UNDER, BG_UNDER, "UNDERLOAD",   "TRAFFIC_STEER SON action\nRecommend handover to enable DTX"),
        (C_SINR,  BG_SINR,  "SINR LOW",    "CRITICAL alert\nRequest PCI re-optimisation from Planning API"),
        (C_PWR,   BG_PWR,   "POWER WASTE", "WARNING alert\nDTX_RECOMMEND: estimated 35% power saving"),
    ]
    for i, (fg, bg, name, action) in enumerate(son_actions):
        x = 0.25 + i * 3.27
        y = 4.05
        rect(s, x, y, 3.1, 0.42, fg, name, fs=10, bold=True, fc=WHITE)
        box(s, x, y + 0.42, 3.1, 1.0, bg, action, fs=9, fc=fg,
            align=PP_ALIGN.LEFT, line_color=fg)

    lbl(s, 0.25, 6.1, 12.8, 0.35,
        "SON actions are also written to InfluxDB 'son_actions' measurement for audit and Grafana dashboard visibility",
        fs=9, fc=DGRAY, italic=True)

    lbl(s, 0.25, 6.5, 12.8, 0.35,
        "NORMAL cells: no action  —  cycle logs: "
        "\"Cycle N | cells=30 | normal=24 | overload=4 | underload=1 | sinr_low=1 | pwr_waste=0\"",
        fs=9, fc=DGRAY)


# ─────────────────────────────────────────────────────────────────────────────
# Main
# ─────────────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    slide_title(prs)
    slide_problem(prs)
    slide_input(prs)
    slide_norm(prs)
    slide_gates(prs)
    slide_bidir(prs)
    slide_arch(prs)
    slide_classes(prs)
    slide_training(prs)
    slide_pipeline(prs)
    prs.save(OUT)
    print(f"Saved -> {OUT}")
    print(f"10 slides: Title, Problem, Input, Normalisation, Gates, "
          f"Bidirectional, Architecture, Classes, Training, Pipeline")


if __name__ == "__main__":
    main()
