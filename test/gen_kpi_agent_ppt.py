"""
Generate KPI_Agent_KPIs.pptx — all KPIs used/collected by the KPI Agent,
with simple plain-English definitions and real Malleswaram values.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "KPI_Agent_KPIs.pptx")

NAVY   = RGBColor(0x0A, 0x29, 0x55)
BLUE   = RGBColor(0x17, 0x5C, 0xA6)
TEAL   = RGBColor(0x00, 0x7E, 0x8A)
GREEN  = RGBColor(0x1B, 0x6F, 0x42)
ORANGE = RGBColor(0xC4, 0x50, 0x08)
PURPLE = RGBColor(0x5A, 0x23, 0x8C)
RED    = RGBColor(0xC0, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF2, 0xF5)
DGRAY  = RGBColor(0x55, 0x55, 0x55)
MGRAY  = RGBColor(0xBB, 0xBB, 0xBB)

BG_NORMAL = RGBColor(0xC6, 0xEF, 0xCE)
BG_OVER   = RGBColor(0xFF, 0xEB, 0x9C)
BG_SINR   = RGBColor(0xFF, 0xC7, 0xCE)
BG_UNDER  = RGBColor(0xBD, 0xD7, 0xEE)
BG_PWR    = RGBColor(0xE2, 0xEF, 0xDA)
BG_GRAY   = RGBColor(0xED, 0xED, 0xED)


def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill, text="", fs=10, bold=False,
         fc=WHITE, align=PP_ALIGN.CENTER):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill
    if text:
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = fc
    return sh

def box(slide, x, y, w, h, fill, text="", fs=10, bold=False,
        fc=WHITE, align=PP_ALIGN.CENTER, lc=None):
    sh = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = lc or fill; sh.line.width = Pt(1.5)
    if text:
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = fc
    return sh

def lbl(slide, x, y, w, h, text, fs=9, fc=NAVY, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold
    r.font.color.rgb = fc; r.font.italic = italic
    return tb

def hdr(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 0.7, NAVY, title, fs=20, bold=True, fc=WHITE)
    if subtitle:
        rect(slide, 0, 0.7, 13.33, 0.30, BLUE, subtitle, fs=9, fc=WHITE)

def pg(slide, n):
    lbl(slide, 12.88, 7.18, 0.32, 0.25, str(n), fs=8, fc=MGRAY,
        align=PP_ALIGN.RIGHT)

def arr(slide, x1, y1, x2, y2, color=NAVY, w=Pt(2)):
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    cx = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cx.line.color.rgb = color; cx.line.width = w
    ln = cx.line._get_or_add_ln()
    etree.SubElement(ln, qn('a:tailEnd'), attrib={'type': 'arrow'})
    return cx


# helper: draw a KPI card  (big name, simple def, range, normal, danger, example)
def kpi_card(s, x, y, w, accent, number,
             name, field, simple_def, unit, normal, danger, ex_normal, ex_danger):
    # accent bar with number
    rect(s, x, y, 0.48, 1.95, accent, str(number), fs=18, bold=True)
    # name + field
    rect(s, x+0.5, y, w-0.5, 0.48, accent, name, fs=12, bold=True)
    lbl(s,  x+0.5, y+0.5, w-0.55, 0.32, field, fs=8, fc=DGRAY, italic=True)
    # simple def
    box(s,  x+0.5, y+0.85, w-0.5, 0.55, LGRAY, simple_def, fs=9, fc=NAVY,
        lc=accent, align=PP_ALIGN.LEFT)
    # normal / danger row
    rect(s, x+0.5, y+1.45, (w-0.5)/2, 0.45, BG_NORMAL,
         "Normal: " + normal, fs=8, bold=True, fc=GREEN)
    rect(s, x+0.5+(w-0.5)/2, y+1.45, (w-0.5)/2, 0.45, BG_SINR,
         "Danger: " + danger, fs=8, bold=True, fc=RED)
    # examples
    rect(s, x+0.5, y+1.95, (w-0.5)/2, 0.38, BG_UNDER,
         "Typical: " + ex_normal, fs=8, fc=BLUE)
    rect(s, x+0.5+(w-0.5)/2, y+1.95, (w-0.5)/2, 0.38, BG_OVER,
         "Overload: " + ex_danger, fs=8, fc=ORANGE)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
def s1(prs):
    s = blank(prs)
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    rect(s, 0, 3.0, 13.33, 0.07, TEAL)

    lbl(s, 0.8, 0.55, 11.7, 1.0,  "KPI Agent", fs=50, fc=WHITE, bold=True,
        align=PP_ALIGN.CENTER)
    lbl(s, 0.8, 1.6,  11.7, 0.65, "All KPIs Used — With Simple Definitions",
        fs=22, fc=RGBColor(0x90, 0xC8, 0xFF), align=PP_ALIGN.CENTER)
    lbl(s, 0.8, 2.38, 11.7, 0.5,
        "Malleswaram Network  ·  30 cells  ·  3 DUs  ·  1 CU  ·  18,400 peak UEs",
        fs=12, fc=WHITE, italic=True, align=PP_ALIGN.CENTER)

    counts = [
        (TEAL,   "9",  "LSTM Input KPIs\n(fed to AI model)"),
        (ORANGE, "5",  "Extra KPIs stored\n(InfluxDB only)"),
        (PURPLE, "7",  "DU-level KPIs\n(infrastructure health)"),
        (GREEN,  "5",  "Health States\n(LSTM output classes)"),
    ]
    for i, (c, big, txt) in enumerate(counts):
        bx = 0.8 + i * 3.0
        box(s, bx, 3.22, 2.65, 1.6, c, lc=c)
        lbl(s, bx+0.1, 3.28, 2.45, 0.75, big,  fs=36, fc=WHITE, bold=True,
            align=PP_ALIGN.CENTER)
        lbl(s, bx+0.1, 4.05, 2.45, 0.65, txt,  fs=10, fc=WHITE,
            align=PP_ALIGN.CENTER)

    lbl(s, 0.8, 5.05, 11.7, 0.38,
        "Source files:  agents/kpi_agent/model.py  |  kpi_agent.py  |  "
        "dev-env/simulators/du/du_simulator.py",
        fs=9, fc=RGBColor(0x70, 0xA0, 0xC8), italic=True, align=PP_ALIGN.CENTER)

    # 5 health states at bottom
    states = [("NORMAL", BG_NORMAL, GREEN), ("OVERLOAD", BG_OVER, ORANGE),
              ("UNDERLOAD", BG_UNDER, BLUE), ("SINR LOW", BG_SINR, RED),
              ("POWER WASTE", BG_PWR, PURPLE)]
    for i, (n, bg, fg) in enumerate(states):
        box(s, 0.4 + i*2.58, 5.62, 2.4, 0.55, bg, n, fs=11, bold=True, fc=fg, lc=fg)
    pg(s, 1)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — All 9 LSTM KPIs at a glance
# ─────────────────────────────────────────────────────────────────────────────
def s2(prs):
    s = blank(prs)
    hdr(s, "9 KPIs Fed Into the LSTM Model — At a Glance",
        "model.py  FEATURE_NORM  |  kpi_agent.py  extract_features()  |  All normalised to 0-1 before inference")
    pg(s, 2)

    kpis = [
        ("#", "KPI Field",          "Simple English Name",
         "Range",          "Normal (5G)",  "Danger Threshold", "Linked Alert"),
        ("1", "prb_dl_pct",         "How busy is the cell?",
         "0 – 100 %",      "40 – 70 %",    "> 85 % → OVERLOAD","OVERLOAD / LOAD_BALANCE"),
        ("2", "sinr_db",            "How clean is the signal?",
         "-5 to +30 dB",   "15 – 22 dB",   "< 5 dB → SINR_LOW","SINR_DEGRADATION / PCI_REOPT"),
        ("3", "connected_ues",      "How many phones connected?",
         "0 – 900",        "300 – 700",    "< 15 + high power","POWER_WASTE / DTX_RECOMMEND"),
        ("4", "power_w",            "How much electricity used?",
         "0 – 1200 W",     "400 – 700 W",  "> 500 W + few UEs","POWER_WASTE / DTX_RECOMMEND"),
        ("5", "packet_loss_pct",    "How many data packets dropped?",
         "0 – 5 %",        "< 0.1 %",      "> 0.5 % at overload","OVERLOAD (secondary)"),
        ("6", "dl_throughput_mbps", "How fast is downlink data?",
         "0 – 4000 Mbps",  "1000–2500 Mbps","Low despite high PRB","SINR_LOW (secondary)"),
        ("7", "cqi",                "Channel quality score (link grade)",
         "0 – 15",         "10 – 14",      "< 5 = poor channel","SINR_LOW (secondary)"),
        ("8", "bler_pct",           "What % of data blocks need resending?",
         "0 – 30 %",       "< 2 %",        "> 8 % at overload","OVERLOAD / SINR_LOW"),
        ("9", "latency_ms",         "How long does one data round-trip take?",
         "0 – 500 ms",     "< 20 ms",      "> 40 ms congested","OVERLOAD (secondary)"),
    ]

    col_w = [0.5, 2.1, 3.3, 1.7, 1.7, 2.0, 1.92]
    y0 = 1.08
    x0 = 0.2

    for ri, row in enumerate(kpis):
        yy = y0 + ri * 0.62
        is_hdr = ri == 0
        bg = NAVY if is_hdr else (LGRAY if ri % 2 == 0 else WHITE)
        fc_row = WHITE if is_hdr else NAVY
        xpos = x0
        for ci, (val, cw) in enumerate(zip(row, col_w)):
            fc_cell = WHITE if is_hdr else (
                TEAL if ci == 1 else
                RED  if ci == 5 else
                GREEN if ci == 4 else fc_row)
            rect(s, xpos, yy, cw, 0.6,
                 bg if is_hdr else (LGRAY if ri % 2 == 0 else WHITE),
                 val, fs=8 if ci > 1 else 9,
                 bold=(is_hdr or ci == 0 or ci == 1),
                 fc=fc_cell if is_hdr else (
                     TEAL if ci == 1 else
                     RED if ci == 5 and ri > 0 else NAVY),
                 align=PP_ALIGN.LEFT if ci >= 2 else PP_ALIGN.CENTER)
            xpos += cw

    lbl(s, 0.2, 6.78, 12.9, 0.35,
        "All 9 values are collected from InfluxDB 'cell_kpi' measurement every poll cycle (POLL_INTERVAL_SEC=10s default).  "
        "6 consecutive readings form one LSTM input tensor  (shape 1 x 6 x 9).",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — KPI 1 & 2: PRB DL % and SINR dB (the most important)
# ─────────────────────────────────────────────────────────────────────────────
def s3(prs):
    s = blank(prs)
    hdr(s, "KPI 1 — PRB DL %   and   KPI 2 — SINR dB   (Most Decision-Critical)",
        "These two KPIs drive OVERLOAD and SINR_LOW detection. Together they reveal if a cell is busy vs degraded.")
    pg(s, 3)

    # PRB
    rect(s, 0.25, 1.05, 6.1, 0.52, TEAL, "KPI 1  —  PRB DL %  (prb_dl_pct)", fs=13, bold=True)
    lbl(s, 0.35, 1.65, 5.9, 0.42,
        "Physical Resource Block Downlink Percentage", fs=10, bold=True, fc=NAVY)
    lbl(s, 0.35, 2.1, 5.9, 0.82,
        "Simple definition:  How much of the cell's radio capacity is being used right now?\n"
        "Think of it like seats on a bus: 50% PRB = half the seats taken. 95% = almost full.",
        fs=10, fc=DGRAY)

    prb_rows = [
        ("0 – 20 %",  "Very quiet",       BG_UNDER, BLUE,   "Off-peak hours, late night"),
        ("20 – 70 %", "Normal operation", BG_NORMAL, GREEN,  "Daytime, moderate load"),
        ("70 – 85 %", "Getting busy",     BG_OVER,   ORANGE, "Evening peak, monitor"),
        ("85 – 98 %", "OVERLOAD",         BG_SINR,   RED,    "Action: move cell to lighter DU"),
    ]
    lbl(s, 0.35, 3.0, 5.9, 0.3, "PRB Range", fs=9, bold=True, fc=NAVY)
    for i, (rng, status, bg, fg, note) in enumerate(prb_rows):
        yy = 3.32 + i * 0.58
        rect(s, 0.35, yy, 1.5,  0.52, bg, rng,    fs=9, bold=True, fc=fg)
        rect(s, 1.88, yy, 1.5,  0.52, bg, status, fs=9, bold=True, fc=fg)
        rect(s, 3.41, yy, 2.9,  0.52, LGRAY, note, fs=9, fc=DGRAY,
             align=PP_ALIGN.LEFT)

    lbl(s, 0.35, 5.7, 5.9, 0.38,
        "Rule threshold: > 85% for one reading (fallback) | LSTM: > 85% sustained for 60s",
        fs=9, fc=DGRAY, italic=True)

    # SINR
    rect(s, 6.95, 1.05, 6.1, 0.52, BLUE, "KPI 2  —  SINR dB  (sinr_db)", fs=13, bold=True)
    lbl(s, 7.05, 1.65, 5.9, 0.42,
        "Signal-to-Interference-plus-Noise Ratio", fs=10, bold=True, fc=NAVY)
    lbl(s, 7.05, 2.1, 5.9, 0.82,
        "Simple definition:  How clear is the radio signal versus background noise?\n"
        "Think of it like trying to hear someone in a noisy room: high SINR = quiet room, low = very noisy.",
        fs=10, fc=DGRAY)

    sinr_rows = [
        ("< 0 dB",    "Terrible signal",  BG_SINR,   RED,    "Almost unusable — severe interference"),
        ("0 – 5 dB",  "SINR LOW",         BG_SINR,   RED,    "Action: CRITICAL alert + PCI reopt"),
        ("5 – 15 dB", "Acceptable",       BG_OVER,   ORANGE, "Degraded but functional"),
        ("> 15 dB",   "Good signal",      BG_NORMAL,  GREEN,  "Normal operation — clear channel"),
    ]
    lbl(s, 7.05, 3.0, 5.9, 0.3, "SINR Range", fs=9, bold=True, fc=NAVY)
    for i, (rng, status, bg, fg, note) in enumerate(sinr_rows):
        yy = 3.32 + i * 0.58
        rect(s, 7.05, yy, 1.4,  0.52, bg, rng,    fs=9, bold=True, fc=fg)
        rect(s, 8.48, yy, 1.5,  0.52, bg, status, fs=9, bold=True, fc=fg)
        rect(s, 10.01,yy, 2.95, 0.52, LGRAY, note, fs=9, fc=DGRAY,
             align=PP_ALIGN.LEFT)

    lbl(s, 7.05, 5.7, 5.9, 0.38,
        "Rule threshold: < 5 dB | LSTM: sustained low SINR at moderate PRB = co-channel interference",
        fs=9, fc=DGRAY, italic=True)

    # Bottom: how they interact
    box(s, 0.25, 6.18, 12.8, 0.72, NAVY, lc=TEAL)
    lbl(s, 0.4, 6.22, 12.5, 0.65,
        "Key interaction:   High PRB + Low SINR  →  Overloaded AND interference  (worst case)\n"
        "High PRB + Good SINR → Pure OVERLOAD (easy to fix — move cell)   "
        "Low PRB + Low SINR → Pure interference (not fixed by moving cell)",
        fs=9, fc=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — KPI 3 & 4: Connected UEs and Power W
# ─────────────────────────────────────────────────────────────────────────────
def s4(prs):
    s = blank(prs)
    hdr(s, "KPI 3 — Connected UEs   and   KPI 4 — Power W",
        "Together these two detect POWER_WASTE: a tower drawing full electricity with almost no users connected.")
    pg(s, 4)

    # UEs
    rect(s, 0.25, 1.05, 6.1, 0.52, ORANGE, "KPI 3  —  Connected UEs  (connected_ues)", fs=12, bold=True)
    lbl(s, 0.35, 1.65, 5.9, 0.38, "How many phones / devices are actively connected to this cell right now?",
        fs=10, bold=True, fc=NAVY)
    lbl(s, 0.35, 2.08, 5.9, 0.52,
        "Simple definition:  The number of active users the tower is serving at this moment.\n"
        "5G NR cell: up to 900 UEs.  4G LTE cell: up to 300 UEs.",
        fs=10, fc=DGRAY)

    ue_rows = [
        ("< 15 UEs",    "Near empty",       BG_PWR,    PURPLE, "POWER_WASTE candidate if power still high"),
        ("15 – 200",    "Low utilisation",  BG_UNDER,  BLUE,   "UNDERLOAD — sleep candidate"),
        ("200 – 700",   "Normal load",      BG_NORMAL, GREEN,  "Healthy operation"),
        ("> 700 UEs",   "Heavy load",       BG_OVER,   ORANGE, "Monitor PRB — approaching OVERLOAD"),
    ]
    for i, (rng, status, bg, fg, note) in enumerate(ue_rows):
        yy = 2.7 + i * 0.58
        rect(s, 0.35, yy, 1.5, 0.52, bg, rng,    fs=9, bold=True, fc=fg)
        rect(s, 1.88, yy, 1.5, 0.52, bg, status, fs=9, bold=True, fc=fg)
        rect(s, 3.41, yy, 2.9, 0.52, LGRAY, note, fs=9, fc=DGRAY, align=PP_ALIGN.LEFT)

    lbl(s, 0.35, 5.02, 5.9, 0.65,
        "Real values (Malleswaram 5G cell):\n"
        "2am off-peak: ~45 UEs   |   2pm normal: ~540 UEs   |   7pm peak: ~870 UEs",
        fs=9, fc=DGRAY, italic=True)

    # Power
    rect(s, 6.95, 1.05, 6.1, 0.52, PURPLE, "KPI 4  —  Power W  (power_w)", fs=12, bold=True)
    lbl(s, 7.05, 1.65, 5.9, 0.38, "How much electrical power is the tower's radio unit consuming right now?",
        fs=10, bold=True, fc=NAVY)
    lbl(s, 7.05, 2.08, 5.9, 0.52,
        "Simple definition:  The electricity bill of the tower. 5G mMIMO is always-on and power-hungry.\n"
        "5G 64T64R: 225–1000W.  4G 4T4R: 50–200W.",
        fs=10, fc=DGRAY)

    pw_rows = [
        ("< 300 W",     "Low power",        BG_NORMAL, GREEN,  "4G cell or 5G in low-load state"),
        ("300–700 W",   "Normal 5G",        BG_NORMAL, GREEN,  "5G cell at 30–70% load"),
        ("700–950 W",   "High 5G",          BG_OVER,   ORANGE, "5G cell near full load"),
        ("> 950 W",     "Near peak",        BG_SINR,   RED,    "Full load — check if OVERLOAD"),
    ]
    for i, (rng, status, bg, fg, note) in enumerate(pw_rows):
        yy = 2.7 + i * 0.58
        rect(s, 7.05, yy, 1.5, 0.52, bg, rng,    fs=9, bold=True, fc=fg)
        rect(s, 8.58, yy, 1.5, 0.52, bg, status, fs=9, bold=True, fc=fg)
        rect(s, 10.11,yy, 2.9, 0.52, LGRAY, note, fs=9, fc=DGRAY, align=PP_ALIGN.LEFT)

    lbl(s, 7.05, 5.02, 5.9, 0.65,
        "Real values (Malleswaram 5G n78, idle_power=237W, tx_power=950W):\n"
        "2am idle: ~250W   |   2pm normal: ~650W   |   7pm peak: ~940W",
        fs=9, fc=DGRAY, italic=True)

    # Combined POWER_WASTE rule
    box(s, 0.25, 5.82, 12.8, 1.0, BG_PWR, lc=PURPLE)
    lbl(s, 0.4, 5.88, 12.5, 0.88,
        "POWER_WASTE is detected when  BOTH conditions are true at the same time:\n"
        "Power W > 500 W   AND   Connected UEs < 15\n"
        "Meaning: the tower is drawing near-full electricity but serving almost nobody — wasted energy.",
        fs=10, bold=False, fc=PURPLE)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — KPI 5, 6: Packet Loss % and Throughput Mbps
# ─────────────────────────────────────────────────────────────────────────────
def s5(prs):
    s = blank(prs)
    hdr(s, "KPI 5 — Packet Loss %   and   KPI 6 — DL Throughput Mbps",
        "Data quality and speed. These confirm what PRB and SINR suggest — if data is actually flowing well.")
    pg(s, 5)

    # Packet Loss
    rect(s, 0.25, 1.05, 6.1, 0.52, RED, "KPI 5  —  Packet Loss %  (packet_loss_pct)", fs=12, bold=True)
    lbl(s, 0.35, 1.65, 5.9, 0.38,
        "What percentage of data packets sent by the network are being dropped (lost)?",
        fs=10, bold=True, fc=NAVY)
    lbl(s, 0.35, 2.08, 5.9, 0.65,
        "Simple definition:  Imagine sending 1000 letters. If 5 never arrive, packet loss = 0.5%.\n"
        "For video calls and downloads, even 1% loss causes buffering and poor quality.",
        fs=10, fc=DGRAY)

    pl_rows = [
        ("0.00 – 0.05 %", "Excellent",  BG_NORMAL, GREEN,  "Healthy network, no congestion"),
        ("0.05 – 0.25 %", "Acceptable", BG_OVER,   ORANGE, "Watch — mild congestion building"),
        ("0.25 – 1.0 %",  "Degraded",   BG_SINR,   RED,    "Cell getting congested or interference"),
        ("> 1.0 %",       "Severe",     BG_SINR,   RED,    "SINR_LOW or OVERLOAD — act now"),
    ]
    for i, (rng, status, bg, fg, note) in enumerate(pl_rows):
        yy = 2.82 + i * 0.58
        rect(s, 0.35, yy, 1.7, 0.52, bg, rng,    fs=9, bold=True, fc=fg)
        rect(s, 2.08, yy, 1.4, 0.52, bg, status, fs=9, bold=True, fc=fg)
        rect(s, 3.51, yy, 2.8, 0.52, LGRAY, note, fs=9, fc=DGRAY, align=PP_ALIGN.LEFT)

    lbl(s, 0.35, 5.22, 5.9, 0.55,
        "How it's computed in simulator:\n"
        "pkt_loss = max(0, (load - 0.75) x 2.5)   — zero below 75% load, rises sharply above",
        fs=9, fc=DGRAY, italic=True)

    # Throughput
    rect(s, 6.95, 1.05, 6.1, 0.52, TEAL, "KPI 6  —  DL Throughput Mbps  (dl_throughput_mbps)", fs=12, bold=True)
    lbl(s, 7.05, 1.65, 5.9, 0.38,
        "How fast is actual downlink data flowing to users right now (megabits per second)?",
        fs=10, bold=True, fc=NAVY)
    lbl(s, 7.05, 2.08, 5.9, 0.65,
        "Simple definition:  The internet speed the tower is delivering.\n"
        "5G NR peak: 3800 Mbps.  4G LTE peak: 150 Mbps.  Affected by both load and signal quality.",
        fs=10, fc=DGRAY)

    tput_rows = [
        ("< 200 Mbps",   "Very low",      BG_SINR,   RED,    "4G cell OR interference — check SINR"),
        ("200 – 1000",   "Moderate",      BG_OVER,   ORANGE, "5G at low load or 4G at peak"),
        ("1000 – 2500",  "Good",          BG_NORMAL,  GREEN,  "5G normal daytime operation"),
        ("> 2500 Mbps",  "Excellent",     BG_NORMAL,  GREEN,  "5G peak performance — near full PRB"),
    ]
    for i, (rng, status, bg, fg, note) in enumerate(tput_rows):
        yy = 2.82 + i * 0.58
        rect(s, 7.05, yy, 1.6,  0.52, bg, rng,    fs=9, bold=True, fc=fg)
        rect(s, 8.68, yy, 1.4,  0.52, bg, status, fs=9, bold=True, fc=fg)
        rect(s, 10.11,yy, 2.9,  0.52, LGRAY, note, fs=9, fc=DGRAY, align=PP_ALIGN.LEFT)

    lbl(s, 7.05, 5.22, 5.9, 0.55,
        "How it's computed: throughput = (prb_dl_pct / 100) x peak_dl_mbps x jitter\n"
        "Key signal: Low throughput DESPITE high PRB = interference stealing efficiency",
        fs=9, fc=DGRAY, italic=True)

    box(s, 0.25, 5.92, 12.8, 0.82, LGRAY, lc=NAVY)
    lbl(s, 0.4, 5.98, 12.5, 0.7,
        "LSTM insight:   Packet loss rising + throughput falling while PRB stays high "
        "= congestion (OVERLOAD).   Low throughput with low PRB "
        "= interference (SINR_LOW).   These patterns are invisible to single-threshold rules.",
        fs=10, fc=NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — KPI 7, 8, 9: CQI, BLER, Latency
# ─────────────────────────────────────────────────────────────────────────────
def s6(prs):
    s = blank(prs)
    hdr(s, "KPI 7 — CQI   |   KPI 8 — BLER %   |   KPI 9 — Latency ms",
        "Channel health indicators. Added in the latest model update. Richer signal for interference detection.")
    pg(s, 6)

    panels = [
        (PURPLE, "KPI 7  —  CQI  (cqi)",
         "Channel Quality Indicator",
         "Simple definition:  A score from 0 to 15 that tells the tower how good the "
         "connection to each user is right now. The tower uses this to choose how fast to send data.\n"
         "Think of it like a Wi-Fi bars indicator — 15 bars = perfect, 0 bars = terrible.",
         [("0 – 4",  "Very poor",  BG_SINR,   RED,    "Interference or severe fading"),
          ("5 – 7",  "Poor",       BG_OVER,   ORANGE, "Degraded — check SINR"),
          ("8 – 11", "Moderate",   BG_NORMAL,  GREEN,  "Normal loaded conditions"),
          ("12 – 15","Excellent",  BG_NORMAL,  GREEN,  "Clean channel, idle or low interference")],
         "Computed from SINR: CQI = (sinr_db + 5) / 2.5  — drops fast when signal degrades",
         0.25),

        (ORANGE, "KPI 8  —  BLER %  (bler_pct)",
         "Block Error Rate",
         "Simple definition:  What percentage of data blocks (chunks of data) have errors "
         "and must be resent? High BLER wastes capacity — the same data is sent multiple times.\n"
         "Like resending a corrupted text message — wastes time and reduces effective speed.",
         [("< 2 %",   "Excellent", BG_NORMAL, GREEN,  "Target range — clean delivery"),
          ("2 – 5 %", "Acceptable",BG_OVER,   ORANGE, "Some retransmissions"),
          ("5 – 10 %","Elevated",  BG_OVER,   ORANGE, "Congestion or signal issues"),
          ("> 10 %",  "High",      BG_SINR,   RED,    "SINR_LOW or severe OVERLOAD")],
         "Computed: bler = (load - 0.75) x 15 + (10 - cqi) x 0.5  — rises sharply near capacity",
         4.44),

        (TEAL, "KPI 9  —  Latency ms  (latency_ms)",
         "Cell Round-Trip Latency",
         "Simple definition:  How long does it take for a single data packet to travel "
         "from the tower to a phone and back? Lower is always better.\n"
         "Like ping time in online gaming — 10ms = fast, 200ms = laggy.",
         [("< 15 ms",  "Excellent", BG_NORMAL, GREEN,  "5G NR normal — ultra-low latency"),
          ("15–30 ms", "Normal",    BG_NORMAL, GREEN,  "Healthy loaded operation"),
          ("30–60 ms", "Elevated",  BG_OVER,   ORANGE, "Congestion building up"),
          ("> 60 ms",  "High",      BG_SINR,   RED,    "OVERLOAD — too many users")],
         "Computed: latency = 8 + load x 25 + max(0, 5 - sinr) x 2  — rises with load and SINR drop",
         8.63),
    ]

    for (clr, title, subtitle, defn, rows, formula, x) in panels:
        w = 4.1
        rect(s, x, 1.05, w, 0.45, clr, title, fs=11, bold=True)
        lbl(s,  x+0.05, 1.55, w-0.1, 0.28, subtitle, fs=9, bold=True, fc=NAVY)
        lbl(s,  x+0.05, 1.85, w-0.1, 0.95, defn, fs=8, fc=DGRAY)
        for i, (rng, status, bg, fg, note) in enumerate(rows):
            yy = 2.88 + i * 0.56
            rect(s, x,       yy, 1.05, 0.5, bg, rng,    fs=8, bold=True, fc=fg)
            rect(s, x+1.07,  yy, 1.05, 0.5, bg, status, fs=8, bold=True, fc=fg)
            rect(s, x+2.14,  yy, 1.92, 0.5, LGRAY, note, fs=7, fc=DGRAY,
                 align=PP_ALIGN.LEFT)
        lbl(s, x+0.05, 5.18, w-0.1, 0.48,
            formula, fs=8, fc=DGRAY, italic=True)

    box(s, 0.25, 5.82, 12.8, 0.92, LGRAY, lc=TEAL)
    lbl(s, 0.4, 5.88, 12.5, 0.8,
        "These 3 KPIs were ADDED in the latest model update (git pull — train.py update, N_FEATURES: 6 -> 9).\n"
        "They give the LSTM a richer picture of interference vs congestion, "
        "especially for distinguishing SINR_LOW from OVERLOAD at the boundary.",
        fs=10, fc=NAVY)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 7 — 5 Extra KPIs Stored (not LSTM input)
# ─────────────────────────────────────────────────────────────────────────────
def s7(prs):
    s = blank(prs)
    hdr(s, "5 Extra KPIs — Stored in InfluxDB but NOT Fed to the LSTM",
        "du_simulator.py pushes these to 'cell_kpi' measurement. Used by Grafana dashboards and Orchestrator queries.")
    pg(s, 7)

    extras = [
        (BLUE,   "UL Throughput Mbps\nul_throughput_mbps",
         "Uplink speed — how fast users upload data TO the tower.\n"
         "Typically 20-25% of DL throughput (asymmetric network design).",
         "0 – 1000 Mbps",
         "Not fed to LSTM — upload problems are rare and usually mirror DL issues.",
         "Grafana: UE Analytics panel. Helps detect asymmetric congestion."),

        (TEAL,   "RSRP dBm\nrsrp_dbm",
         "Reference Signal Received Power — absolute received signal strength at the UE.\n"
         "Stronger (less negative) = UE is physically closer to the tower.",
         "-120 to -60 dBm",
         "Not fed to LSTM — SINR already captures signal quality more completely.",
         "Grafana: Cell KPI panel. Used for handover decisions and coverage analysis."),

        (ORANGE, "PRB UL %\nprb_ul_pct",
         "Uplink Physical Resource Block usage — how much upload capacity is occupied.\n"
         "Usually much lower than DL PRB in typical consumer traffic patterns.",
         "0 – 95 %",
         "Not fed to LSTM — highly correlated with DL PRB so adds little new information.",
         "Grafana: Cell KPI panel. Useful for IoT/mMTC cells that are uplink-heavy."),

        (GREEN,  "HO Success Rate\nho_success_rate",
         "Handover Success Rate — what fraction of phone-to-phone-tower switches complete without dropping the call.\n"
         "Target: > 96%. Drops indicate signalling problems or coverage holes.",
         "0 – 100 %",
         "Not fed to LSTM — captures rare events that are hard to model in 60s windows.",
         "Grafana: Network Overview. Alerts if drops below 96% threshold."),

        (PURPLE, "Coverage Radius m\ncoverage_radius_m",
         "Estimated RF coverage radius computed from COST-231-Hata path loss model.\n"
         "Larger radius = lower frequency band (n28 700 MHz vs n78 3500 MHz).",
         "200 – 2000 m",
         "Not fed to LSTM — static per cell, does not change at runtime. Config-time only.",
         "Used by Planning API for capacity planning and DU assignment."),
    ]

    for i, (clr, name, defn, rng, why_not, where) in enumerate(extras):
        y = 1.08 + i * 1.24
        rect(s, 0.25, y, 2.8, 1.1, clr, name, fs=9, bold=True)
        lbl(s,  3.1,  y, 5.6, 0.5, defn, fs=9, fc=DGRAY)
        rect(s, 3.1,  y+0.52, 1.5, 0.55, LGRAY, rng, fs=8, fc=NAVY)
        rect(s, 4.63, y+0.52, 4.05, 0.55, BG_UNDER, why_not, fs=8, fc=BLUE,
             align=PP_ALIGN.LEFT)
        rect(s, 8.72, y, 4.4, 1.1, LGRAY, where, fs=9, fc=DGRAY,
             align=PP_ALIGN.LEFT)

    for x, w, h in [(0.25,2.8,"Extra KPI"), (3.1,5.6,"Definition"),
                    (8.72,4.4,"Where it is used")]:
        rect(s, x, 0.76, w, 0.3, NAVY, h, fs=8, bold=True, fc=WHITE)

    lbl(s, 0.25, 7.12, 12.8, 0.3,
        "All 14 cell_kpi fields flow through InfluxDB. The KPI agent reads only the 9 it needs. "
        "Grafana dashboards visualise all 14.",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 8 — DU-Level KPIs
# ─────────────────────────────────────────────────────────────────────────────
def s8(prs):
    s = blank(prs)
    hdr(s, "7 DU-Level KPIs — Infrastructure Health of Each Distributed Unit",
        "Written to 'du_kpi' measurement in InfluxDB every tick. Covers compute, memory, fronthaul, and signalling.")
    pg(s, 8)

    du_kpis = [
        (TEAL,   "active_ues",             "Total UEs",
         "Total number of phones connected across all cells on this DU right now.",
         "DU-MLS-1 (12 cells): 0 – 10,800",
         "Sum of connected_ues per cell. High = DU is busy."),
        (TEAL,   "cell_count",             "Active Cells",
         "How many cells are currently assigned to and active on this DU.",
         "8 – 12 cells",
         "Changes when KPI agent moves cells between DUs."),
        (ORANGE, "cpu_pct",                "CPU Usage %",
         "How much of the DU's processing power is used. L1/L2 baseband processing is CPU-heavy.",
         "0 – 100 %",
         "cpu = 20 + load x 62. High at peak hours or after cell additions."),
        (ORANGE, "memory_pct",             "Memory Usage %",
         "How much RAM the DU is using. Stores UE contexts, buffers, scheduling queues.",
         "0 – 100 %",
         "memory = 30 + load x 45. Spikes if many cells assigned."),
        (BLUE,   "fronthaul_latency_us",   "Fronthaul Latency",
         "Delay in microseconds on the eCPRI/CPRI link between DU and RU (radio unit on the tower).",
         "50 – 200 microseconds",
         "Must stay below 100 us for 5G NR timing requirements."),
        (BLUE,   "processing_delay_ms",    "Processing Delay",
         "Time the DU takes to process one L1/L2 scheduling cycle.",
         "0.1 – 0.9 ms",
         "Increases under high CPU load. > 1 ms causes scheduling jitter."),
        (PURPLE, "f1_msg_per_sec",         "F1 Messages / sec",
         "Rate of F1-AP control messages exchanged between the DU and CU over the F1 interface.",
         "0 – 20,000 / sec",
         "f1_rate = active_ues x 0.5–2.0. Peaks during handovers and reconfigs."),
    ]

    col_w = [2.2, 1.4, 3.6, 2.4, 3.2]
    y0 = 1.05
    for ri, (clr, field, name, defn, rng, note) in enumerate(du_kpis):
        yy = y0 + ri * 0.82
        bg = LGRAY if ri % 2 == 0 else WHITE
        rect(s, 0.25, yy, 0.35, 0.76, clr)
        rect(s, 0.63, yy, 2.1,  0.76, bg, field,
             fs=9, bold=True, fc=clr, align=PP_ALIGN.LEFT)
        rect(s, 2.76, yy, 1.5,  0.76, bg, name,
             fs=9, bold=True, fc=NAVY)
        rect(s, 4.29, yy, 4.3,  0.76, bg, defn,
             fs=9, fc=DGRAY, align=PP_ALIGN.LEFT)
        rect(s, 8.62, yy, 1.85, 0.76, bg, rng,
             fs=8, fc=NAVY, align=PP_ALIGN.LEFT)
        rect(s, 10.5, yy, 2.6,  0.76, bg, note,
             fs=8, fc=DGRAY, align=PP_ALIGN.LEFT)

    for x, w, h in [(0.63,2.1,"Field"), (2.76,1.5,"Name"), (4.29,4.3,"Definition"),
                    (8.62,1.85,"Range"), (10.5,2.6,"How computed / note")]:
        rect(s, x, 0.75, w, 0.3, NAVY, h, fs=8, bold=True, fc=WHITE)

    lbl(s, 0.25, 6.82, 12.8, 0.35,
        "DU KPIs are NOT fed to the LSTM. They are monitored via Grafana (du_cu_performance dashboard) "
        "and used by the Orchestrator's query_network tool to report DU health to the LLM.",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 9 — KPI → Health State mapping
# ─────────────────────────────────────────────────────────────────────────────
def s9(prs):
    s = blank(prs)
    hdr(s, "KPI Signatures — What Each Health State Looks Like in the Data",
        "How the LSTM learned to associate KPI patterns with the 5 health states (from train.py class specs)")
    pg(s, 9)

    # Header
    cols = ["Health State", "PRB DL %", "SINR dB", "UEs", "Power W",
            "Pkt Loss %", "Tput Mbps", "CQI", "BLER %", "Latency ms"]
    cw   = [1.85, 1.15, 1.05, 0.85, 1.05, 1.1, 1.1, 0.75, 0.85, 1.15]
    y0   = 1.05
    xpos = 0.25
    for h, w in zip(cols, cw):
        rect(s, xpos, y0, w, 0.4, NAVY, h, fs=8, bold=True, fc=WHITE)
        xpos += w

    # 5G values from train.py _5G_SPECS means
    rows = [
        (GREEN,  BG_NORMAL, "NORMAL (5G)",
         "55%", "20 dB", "350", "520 W", "0.05%", "1400", "11", "1.5%", "12 ms"),
        (ORANGE, BG_OVER,   "OVERLOAD (5G)",
         "94%", "11 dB", "720", "940 W", "0.85%", "3100", "7",  "8.0%", "38 ms"),
        (BLUE,   BG_UNDER,  "UNDERLOAD (5G)",
         "9%",  "24 dB", "20",  "330 W", "0.01%", "190",  "14", "0.3%", "9 ms"),
        (RED,    BG_SINR,   "SINR LOW (5G)",
         "54%", "1 dB",  "290", "580 W", "1.60%", "720",  "3",  "12%",  "45 ms"),
        (PURPLE, BG_PWR,    "POWER WASTE (5G)",
         "13%", "24 dB", "8",   "880 W", "0.01%", "145",  "14", "0.2%", "9 ms"),
    ]

    for ri, (fg, bg, state, *vals) in enumerate(rows):
        yy = y0 + 0.4 + ri * 1.0
        # state label
        rect(s, 0.25, yy, 1.85, 0.95, bg, state, fs=10, bold=True, fc=fg)
        xpos = 0.25 + 1.85
        for vi, (v, w) in enumerate(zip(vals, cw[1:])):
            # highlight danger values
            is_bad = ((ri==1 and vi in [0,1,4,6]) or
                      (ri==3 and vi in [1,4,6]) or
                      (ri==4 and vi in [3]))
            cell_bg = BG_SINR if is_bad else bg
            cell_fc = RED     if is_bad else fg
            rect(s, xpos, yy, w, 0.95, cell_bg, v,
                 fs=10, bold=is_bad, fc=cell_fc)
            xpos += w

    lbl(s, 0.25, 6.08, 12.8, 0.7,
        "Red cells = the KPI values that are most diagnostic for that health state.\n"
        "The LSTM learns to recognise these patterns across 6 consecutive timesteps — "
        "not just the snapshot values shown here, but the TREND leading up to them.",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 10 — Summary: Full KPI Inventory
# ─────────────────────────────────────────────────────────────────────────────
def s10(prs):
    s = blank(prs)
    hdr(s, "Summary — Full KPI Inventory of the KPI Agent",
        "21 total KPIs collected and used across the Malleswaram 4G/5G deployment")
    pg(s, 10)

    # Three columns
    groups = [
        (TEAL, "9 LSTM Input KPIs",
         "Fed to AI model every cycle\n(model.py FEATURE_NORM)",
         ["1.  prb_dl_pct      — Cell busy-ness %",
          "2.  sinr_db         — Signal clarity (dB)",
          "3.  connected_ues   — Active phones",
          "4.  power_w         — Electricity used (W)",
          "5.  packet_loss_pct — Dropped data %",
          "6.  dl_throughput_mbps — Downlink speed",
          "7.  cqi             — Channel quality score",
          "8.  bler_pct        — Retransmission rate %",
          "9.  latency_ms      — Round-trip time (ms)"]),

        (ORANGE, "5 Extra Cell KPIs",
         "Stored in InfluxDB / Grafana\n(not LSTM input)",
         ["10. ul_throughput_mbps — Upload speed",
          "11. rsrp_dbm          — Signal strength (dBm)",
          "12. prb_ul_pct        — Upload capacity used %",
          "13. ho_success_rate   — Handover success %",
          "14. coverage_radius_m — RF coverage area (m)"]),

        (PURPLE, "7 DU Infrastructure KPIs",
         "Per Distributed Unit health\n('du_kpi' measurement)",
         ["15. active_ues           — Total connected UEs",
          "16. cell_count           — Cells on this DU",
          "17. cpu_pct              — CPU usage %",
          "18. memory_pct           — RAM usage %",
          "19. fronthaul_latency_us — DU-to-RU delay",
          "20. processing_delay_ms  — Scheduling delay",
          "21. f1_msg_per_sec       — F1 interface rate"]),
    ]

    for i, (clr, title, subtitle, items) in enumerate(groups):
        x = 0.25 + i * 4.35
        rect(s, x, 1.05, 4.1, 0.52, clr, title, fs=13, bold=True)
        lbl(s, x+0.1, 1.62, 3.9, 0.48, subtitle, fs=9, fc=DGRAY, italic=True)
        for j, item in enumerate(items):
            yy = 2.15 + j * 0.52
            bg = LGRAY if j % 2 == 0 else WHITE
            rect(s, x, yy, 4.1, 0.5, bg, item, fs=9, fc=NAVY,
                 align=PP_ALIGN.LEFT)

    # Bottom summary
    box(s, 0.25, 7.05, 12.8, 0.3, NAVY, lc=NAVY)
    lbl(s, 0.4, 7.07, 12.5, 0.25,
        "KPI Agent uses 9 of 21 KPIs for AI classification  ·  "
        "All 21 flow through InfluxDB  ·  Grafana shows all via 5 dashboards  ·  "
        "Orchestrator LLM queries any of them via query_network / query_cell tools",
        fs=9, fc=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    s1(prs); s2(prs); s3(prs); s4(prs); s5(prs)
    s6(prs); s7(prs); s8(prs); s9(prs); s10(prs)
    prs.save(OUT)
    print(f"Saved -> {OUT}")
    print("10 slides: Title, 9 KPIs at a glance, PRB+SINR deep-dive, "
          "UEs+Power, Pkt Loss+Tput, CQI+BLER+Latency, "
          "Extra KPIs, DU KPIs, Health State signatures, Full inventory")

if __name__ == "__main__":
    main()
