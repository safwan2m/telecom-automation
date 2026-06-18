"""
Generate LSTM_Story.pptx — simple story-style explanation of the KPI LSTM classifier.
No jargon. Analogy-first. Real data from the Malleswaram deployment.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "LSTM_Story.pptx")

# ── colours ───────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x29, 0x55)
BLUE   = RGBColor(0x17, 0x5C, 0xA6)
TEAL   = RGBColor(0x00, 0x7E, 0x8A)
GREEN  = RGBColor(0x1B, 0x6F, 0x42)
ORANGE = RGBColor(0xC4, 0x50, 0x08)
PURPLE = RGBColor(0x5A, 0x23, 0x8C)
RED    = RGBColor(0xC0, 0x00, 0x00)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY  = RGBColor(0xF0, 0xF2, 0xF5)
DGRAY  = RGBColor(0x50, 0x50, 0x50)
MGRAY  = RGBColor(0xCC, 0xCC, 0xCC)

BG_NORMAL = RGBColor(0xC6, 0xEF, 0xCE)
BG_OVER   = RGBColor(0xFF, 0xEB, 0x9C)
BG_UNDER  = RGBColor(0xBD, 0xD7, 0xEE)
BG_SINR   = RGBColor(0xFF, 0xC7, 0xCE)
BG_PWR    = RGBColor(0xE2, 0xEF, 0xDA)


def new_prs():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill, text="", fs=10, bold=False,
         fc=WHITE, align=PP_ALIGN.CENTER, wrap=True):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = fill
    if text:
        tf = sh.text_frame; tf.word_wrap = wrap
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = fc
    return sh

def box(slide, x, y, w, h, fill, text="", fs=10, bold=False,
        fc=WHITE, align=PP_ALIGN.CENTER, lc=None):
    sh = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = lc or fill; sh.line.width = Pt(1.5)
    if text:
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(fs); r.font.bold = bold; r.font.color.rgb = fc
    return sh

def lbl(slide, x, y, w, h, text, fs=9, fc=NAVY, bold=False,
        align=PP_ALIGN.CENTER, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.bold = bold
    r.font.color.rgb = fc; r.font.italic = italic
    return tb

def arr(slide, x1, y1, x2, y2, color=NAVY, w=Pt(2.5)):
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    cx = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cx.line.color.rgb = color; cx.line.width = w
    ln = cx.line._get_or_add_ln()
    etree.SubElement(ln, qn('a:tailEnd'), attrib={'type': 'arrow'})
    return cx

def hdr(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 0.7, NAVY, title, fs=21, bold=True, fc=WHITE)
    if subtitle:
        rect(slide, 0, 0.7, 13.33, 0.32, BLUE, subtitle, fs=10, fc=WHITE)

def pg(slide, n):
    lbl(slide, 12.85, 7.15, 0.35, 0.28, str(n), fs=8, fc=MGRAY)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 1 — Title
# ─────────────────────────────────────────────────────────────────────────────
def s1_title(prs):
    s = blank(prs)
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    rect(s, 0, 3.55, 13.33, 0.07, TEAL)

    lbl(s, 0.8, 0.7,  11.7, 1.1,  "What is LSTM?",
        fs=48, fc=WHITE, bold=True)
    lbl(s, 0.8, 1.85, 11.7, 0.65, "Simple Explanation with Real Data",
        fs=22, fc=RGBColor(0x90, 0xC8, 0xFF))
    lbl(s, 0.8, 2.6,  11.7, 0.7,
        "How an AI watches 30 mobile cells every 10 seconds\n"
        "and keeps the Malleswaram 4G/5G network healthy",
        fs=13, fc=WHITE, italic=True)

    # 3 big facts
    facts = [
        (TEAL,   "30",  "Cells monitored"),
        (ORANGE, "9",   "KPIs watched"),
        (PURPLE, "60s", "History per decision"),
    ]
    for i, (c, big, small) in enumerate(facts):
        bx = 1.5 + i * 3.7
        box(s, bx, 3.8, 3.0, 1.5, c, fc=WHITE, lc=c)
        lbl(s, bx, 3.9, 3.0, 0.8, big,   fs=36, fc=WHITE, bold=True)
        lbl(s, bx, 4.75, 3.0, 0.45, small, fs=11, fc=WHITE)

    lbl(s, 0.8, 5.6, 11.7, 0.45,
        "5 possible diagnoses:   NORMAL     OVERLOAD     UNDERLOAD     SINR LOW     POWER WASTE",
        fs=11, fc=RGBColor(0xA0, 0xD0, 0xFF))
    lbl(s, 0.8, 6.1, 11.7, 0.35,
        "agents/kpi_agent/model.py  |  train.py  |  kpi_agent.py",
        fs=9, fc=RGBColor(0x60, 0x90, 0xC0), italic=True)
    pg(s, 1)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 2 — The Doctor Analogy
# ─────────────────────────────────────────────────────────────────────────────
def s2_analogy(prs):
    s = blank(prs)
    hdr(s, "Think of LSTM Like a Network Doctor",
        "A doctor does not diagnose from one reading — they look at the trend over time. LSTM does the same.")
    pg(s, 2)

    # Left — Doctor
    rect(s, 0.3, 1.1, 5.9, 0.52, TEAL, "DOCTOR checks a patient", fs=13, bold=True)
    rows_doc = [
        ("Checks vitals every few minutes",   "Polls each cell every 10 seconds"),
        ("Reviews last 6 readings on chart",  "Reads last 6 KPI snapshots (60 s)"),
        ("Looks at 9 vital signs",            "Watches 9 network KPIs per cell"),
        ("Diagnoses: healthy / sick / urgent","Classifies: NORMAL / OVERLOAD / SINR_LOW..."),
        ("Prescribes treatment",              "Triggers SON action (move cell, alert, DTX)"),
    ]
    for i, (doc, net) in enumerate(rows_doc):
        yy = 1.68 + i * 0.72
        bg = LGRAY if i % 2 == 0 else WHITE
        rect(s, 0.3, yy, 5.9, 0.65, bg, doc, fs=10, fc=NAVY, align=PP_ALIGN.LEFT)

    # Right — Network
    rect(s, 7.1, 1.1, 5.9, 0.52, BLUE, "LSTM checks a cell", fs=13, bold=True)
    for i, (doc, net) in enumerate(rows_doc):
        yy = 1.68 + i * 0.72
        bg = LGRAY if i % 2 == 0 else WHITE
        rect(s, 7.1, yy, 5.9, 0.65, bg, net, fs=10, fc=NAVY, align=PP_ALIGN.LEFT)

    # VS divider
    lbl(s, 6.25, 3.2, 0.8, 0.6, "=", fs=28, bold=True, fc=NAVY)

    # Key insight box
    box(s, 0.3, 5.42, 12.7, 0.92, NAVY, fc=WHITE, lc=TEAL)
    lbl(s, 0.3, 5.5, 12.7, 0.75,
        "Key Insight:  A single reading of PRB = 86% could be a harmless spike — or the start of a real overload.\n"
        "The LSTM sees 6 readings in a row and knows the difference. Rules cannot.",
        fs=11, bold=False, fc=WHITE, italic=False)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 3 — 9 Things It Watches (KPIs in plain English)
# ─────────────────────────────────────────────────────────────────────────────
def s3_kpis(prs):
    s = blank(prs)
    hdr(s, "9 Things the LSTM Watches — Per Cell, Every 10 Seconds",
        "These are the 9 KPI features fed into the model  (model.py  N_FEATURES = 9)")
    pg(s, 3)

    kpis = [
        (TEAL,   "PRB DL %",
         "How busy is the cell?",
         "0% = idle   100% = completely full",
         ">85% for 60s = OVERLOAD"),
        (TEAL,   "SINR dB",
         "How clear is the signal?",
         "Higher = cleaner signal",
         "<5 dB for 60s = SINR_LOW"),
        (TEAL,   "Connected UEs",
         "How many phones are connected?",
         "5G cell: up to 900 phones",
         "Very few + high power = POWER_WASTE"),
        (BLUE,   "Power W",
         "How much electricity is the tower using?",
         "5G idle: ~250 W   Full load: ~1000 W",
         "High power + few users = waste"),
        (BLUE,   "Packet Loss %",
         "How many data packets are being dropped?",
         "Target: <0.1%  Overload: >0.5%",
         "Rising with PRB = congestion"),
        (BLUE,   "Throughput Mbps",
         "How fast is data actually flowing?",
         "5G peak: 3800 Mbps  4G: 150 Mbps",
         "Low despite high PRB = interference"),
        (PURPLE, "CQI",
         "Channel Quality Indicator (link health score)",
         "0 = terrible   15 = perfect",
         "<5 = poor channel quality"),
        (PURPLE, "BLER %",
         "Block Error Rate — how many data blocks need resending?",
         "Target: <10%   Overload: >8%",
         "High BLER + low SINR = interference"),
        (PURPLE, "Latency ms",
         "How long does one round-trip take?",
         "Normal: <20 ms   Overload: >40 ms",
         "Rising latency = congestion"),
    ]

    col_w = [2.2, 2.9, 2.5, 3.2]
    x0, y0 = 0.25, 1.12
    # Header
    for ci, (h, w) in enumerate(zip(
            ["KPI Name", "Plain English", "Range", "Warning Sign"], col_w)):
        xpos = x0 + sum(col_w[:ci])
        rect(s, xpos, y0, w, 0.38, NAVY, h, fs=9, bold=True, fc=WHITE)

    for ri, (clr, name, plain, rng, warn) in enumerate(kpis):
        yy = y0 + 0.38 + ri * 0.56
        bg = LGRAY if ri % 2 == 0 else WHITE
        vals = [name, plain, rng, warn]
        fcs  = [clr, DGRAY, DGRAY, RED]
        bolds = [True, False, False, False]
        for ci, (v, fc, b, w) in enumerate(zip(vals, fcs, bolds, col_w)):
            xpos = x0 + sum(col_w[:ci])
            rect(s, xpos, yy, w, 0.54, bg, v, fs=9 if ci > 0 else 10,
                 bold=b, fc=fc, align=PP_ALIGN.LEFT)

    lbl(s, 0.25, 6.82, 12.8, 0.35,
        "All 9 values are normalised to 0-1 before feeding into the LSTM  "
        "(so PRB=94% becomes 0.94 and Power=950W becomes 0.79 — same scale)",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 4 — How LSTM Remembers (Simple)
# ─────────────────────────────────────────────────────────────────────────────
def s4_memory(prs):
    s = blank(prs)
    hdr(s, "How the LSTM Remembers — Two Types of Memory",
        "Unlike a simple rule, the LSTM carries a memory that updates at every timestep")
    pg(s, 4)

    # Two memory types
    rect(s, 0.3, 1.05, 6.0, 0.55, TEAL,
         "LONG-TERM MEMORY  (Cell State)", fs=13, bold=True)
    box(s, 0.3, 1.65, 6.0, 2.5, RGBColor(0xE8, 0xF8, 0xFA),
        fc=NAVY, lc=TEAL)
    lbl(s, 0.5, 1.8, 5.6, 2.2,
        "Carries important signals across all 6 timesteps\n\n"
        "Example:  \"PRB has been rising\n"
        "            for the last 5 readings\"\n\n"
        "This signal travels from t-50s all the\n"
        "way to t=now without being lost or forgotten",
        fs=11, fc=DGRAY, align=PP_ALIGN.LEFT)

    rect(s, 7.0, 1.05, 6.0, 0.55, ORANGE,
         "SHORT-TERM MEMORY  (Hidden State)", fs=13, bold=True)
    box(s, 7.0, 1.65, 6.0, 2.5, RGBColor(0xFF, 0xF8, 0xEE),
        fc=NAVY, lc=ORANGE)
    lbl(s, 7.2, 1.8, 5.6, 2.2,
        "What just happened — passed to the next step\n\n"
        "Example:  \"At t-10s the PRB\n"
        "            jumped sharply to 89%\"\n\n"
        "This recent event influences the decision\n"
        "at the very next timestep (t=now)",
        fs=11, fc=DGRAY, align=PP_ALIGN.LEFT)

    # 3 gates simplified
    lbl(s, 0.3, 4.35, 12.7, 0.4,
        "Three gates control both memories at every step:", fs=12, bold=True, fc=NAVY)

    gates = [
        (ORANGE, "FORGET",
         "Throw away old info\nthat no longer matters",
         "\"Last night's low-traffic\n readings — forget them\""),
        (TEAL,   "INPUT",
         "Store something new\nthat just happened",
         "\"PRB just jumped to 89%\n — remember that\""),
        (PURPLE, "OUTPUT",
         "Decide what to pass\nforward to next step",
         "\"Tell next step:\n rising PRB trend\""),
    ]
    for i, (clr, name, desc, ex) in enumerate(gates):
        gx = 0.3 + i * 4.35
        rect(s, gx, 4.82, 4.0, 0.48, clr, name + " GATE", fs=12, bold=True)
        box(s, gx, 5.35, 4.0, 0.88, LGRAY, desc, fs=10, fc=NAVY, lc=clr)
        box(s, gx, 6.28, 4.0, 0.85, WHITE, ex,   fs=10, fc=clr,  lc=clr)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 5 — Reading Forward AND Backward
# ─────────────────────────────────────────────────────────────────────────────
def s5_bidir(prs):
    s = blank(prs)
    hdr(s, "Bidirectional — The LSTM Reads the 60-Second Window Both Ways",
        "Two LSTMs run in parallel: one reads t-50s to t=now, the other reads t=now back to t-50s")
    pg(s, 5)

    # Timeline
    ts   = ["t-50s", "t-40s", "t-30s", "t-20s", "t-10s", "t=now"]
    prbs = [ "60%",   "68%",   "75%",   "82%",   "89%",   "94%"]
    cols = [LGRAY, LGRAY, BG_OVER, BG_OVER, BG_SINR, BG_SINR]
    fcs  = [NAVY,  NAVY,  ORANGE, ORANGE,   RED,     RED]
    for i in range(6):
        x = 0.5 + i * 2.05
        rect(s, x, 1.05, 1.85, 0.52, cols[i], ts[i],   fs=9,  bold=False, fc=fcs[i])
        rect(s, x, 1.62, 1.85, 0.62, cols[i], prbs[i], fs=16, bold=True,  fc=fcs[i])

    # Forward arrow
    arr(s, 0.8, 2.55, 12.2, 2.55, color=TEAL, w=Pt(3))
    rect(s, 0.3, 2.65, 12.7, 0.52, TEAL,
         "Forward LSTM reads left to right:  \"PRB is climbing... 60...68...75...82...89...94\"  → OVERLOAD coming",
         fs=10, bold=True)

    # Backward arrow
    arr(s, 12.2, 3.52, 0.8, 3.52, color=PURPLE, w=Pt(3))
    rect(s, 0.3, 3.6, 12.7, 0.52, PURPLE,
         "Backward LSTM reads right to left: \"t=now is 94%, previous steps confirm it was a real sustained trend\"",
         fs=10, bold=True)

    # Combine box
    box(s, 1.5, 4.35, 10.3, 0.72, NAVY, fc=WHITE, lc=TEAL)
    lbl(s, 1.5, 4.38, 10.3, 0.65,
        "Both readings combined  →  \"6 consecutive rising readings, confirmed from both ends = OVERLOAD\"",
        fs=12, bold=True, fc=WHITE)

    # Why it matters
    lbl(s, 0.3, 5.28, 5.9, 0.38, "Without bidirectional:", fs=11, bold=True, fc=RED)
    box(s, 0.3, 5.68, 5.9, 1.4, BG_SINR, fc=RED, lc=RED)
    lbl(s, 0.5, 5.78, 5.5, 1.2,
        "A one-directional LSTM might see the spike\n"
        "at t=now but not know whether what came\n"
        "before it was a real trend or random noise",
        fs=10, fc=DGRAY, align=PP_ALIGN.LEFT)

    lbl(s, 6.5, 5.28, 6.5, 0.38, "With bidirectional:", fs=11, bold=True, fc=GREEN)
    box(s, 6.5, 5.68, 6.5, 1.4, BG_NORMAL, fc=GREEN, lc=GREEN)
    lbl(s, 6.7, 5.78, 6.1, 1.2,
        "The backward pass confirms:\n"
        "PRB was rising for ALL 6 steps.\n"
        "This is definitely a real overload —\n"
        "act now with high confidence.",
        fs=10, fc=DGRAY, align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 6 — 5 Decisions and What Happens
# ─────────────────────────────────────────────────────────────────────────────
def s6_decisions(prs):
    s = blank(prs)
    hdr(s, "5 Decisions the LSTM Makes — and What Happens Next",
        "After reading 6 timesteps, the model outputs a probability for each class. The highest wins.")
    pg(s, 6)

    classes = [
        (GREEN,  BG_NORMAL, "0  NORMAL",
         "PRB 40-70%  |  SINR good  |  Stable UEs",
         "70% of time",
         "Nothing — network is healthy"),
        (ORANGE, BG_OVER,   "1  OVERLOAD",
         "PRB >85%  |  Packet loss rising  |  Many UEs",
         "15% of time",
         "Auto-move cell to lightest DU\n(cooldown: 30s between moves)"),
        (BLUE,   BG_UNDER,  "2  UNDERLOAD",
         "PRB <20%  |  Very few users  |  Wasted capacity",
         "8% of time",
         "Flag as sleep candidate\nRecommend DTX energy saving"),
        (RED,    BG_SINR,   "3  SINR LOW",
         "SINR <5 dB  |  High errors  |  Poor CQI",
         "5% of time",
         "CRITICAL alert raised\nRequest PCI re-optimisation"),
        (PURPLE, BG_PWR,    "4  POWER WASTE",
         "High power  |  Very few users (<15)  |  Low PRB",
         "2% of time",
         "WARNING alert\nRecommend DTX / sleep mode\nEstimated 35% power saving"),
    ]

    for i, (fg, bg, name, kpis, freq, action) in enumerate(classes):
        y = 1.08 + i * 1.22
        rect(s, 0.3, y, 2.8, 1.08, bg, name, fs=13, bold=True, fc=fg)
        rect(s, 3.15, y, 3.8, 1.08, LGRAY if i%2==0 else WHITE,
             kpis, fs=9, fc=DGRAY, align=PP_ALIGN.LEFT)
        rect(s, 7.0, y, 1.5, 1.08, bg, freq, fs=10, bold=True, fc=fg)
        box(s, 8.55, y, 4.5, 1.08, WHITE, action, fs=10, fc=fg,
            align=PP_ALIGN.LEFT, lc=fg)

    # Column headers
    for x, w, h in [(0.3,2.8,"Class"), (3.15,3.8,"KPI Signature"),
                    (7.0,1.5,"Frequency"), (8.55,4.5,"SON Action Triggered")]:
        rect(s, x, 0.75, w, 0.32, NAVY, h, fs=9, bold=True, fc=WHITE)

    lbl(s, 0.3, 7.12, 12.7, 0.3,
        "MIN_CONFIDENCE = 70% — if the model is less than 70% sure, it logs the finding but does NOT act",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 7 — Proof: Validation Test Results
# ─────────────────────────────────────────────────────────────────────────────
def s7_results(prs):
    s = blank(prs)
    hdr(s, "Proof — Validation Test Results  (30 runs per scenario, seed=42)",
        "We tested the trained LSTM against 7 synthetic scenarios. Here is what it predicted.")
    pg(s, 7)

    # Summary badges
    box(s, 0.3, 1.05, 3.8, 1.1, BG_NORMAL, fc=GREEN, lc=GREEN)
    lbl(s, 0.3, 1.1, 3.8, 0.55, "4 / 7", fs=32, bold=True, fc=GREEN)
    lbl(s, 0.3, 1.65, 3.8, 0.42, "Scenarios PASSED  (100% confidence)", fs=10, fc=GREEN)

    box(s, 4.3, 1.05, 3.8, 1.1, BG_SINR, fc=RED, lc=RED)
    lbl(s, 4.3, 1.1,  3.8, 0.55, "3 / 7", fs=32, bold=True, fc=RED)
    lbl(s, 4.3, 1.65, 3.8, 0.42, "Scenarios FAILED  (expose real gaps)", fs=10, fc=RED)

    box(s, 8.3, 1.05, 4.7, 1.1, BG_UNDER, fc=BLUE, lc=BLUE)
    lbl(s, 8.3, 1.1,  4.7, 0.55, "100%", fs=32, bold=True, fc=BLUE)
    lbl(s, 8.3, 1.65, 4.7, 0.42, "Confidence on every prediction (LSTM never unsure)", fs=10, fc=BLUE)

    # Results table
    results = [
        (True,  "NORMAL — 5G cell at 50% load",
         "NORMAL",    "NORMAL",  "100%", "Correct — quiet cell correctly left alone"),
        (False, "NORMAL — 4G cell at 55% load",
         "NORMAL",    "UNDERLOAD","100%","4G has fewer UEs/lower power — LSTM thinks it is underloaded\n(training data was 5G-heavy)"),
        (True,  "OVERLOAD — 5G cell at 94% load",
         "OVERLOAD",  "OVERLOAD","100%", "Correct — auto-move would be triggered"),
        (False, "OVERLOAD — 4G cell at 94% load",
         "OVERLOAD",  "SINR_LOW","100%","4G throughput (141 Mbps) looks like interference to model\n(training assumed 5G-scale 3100 Mbps at overload)"),
        (True,  "SINR_LOW — interference injected",
         "SINR_LOW",  "SINR_LOW","100%","Correct — CRITICAL alert + PCI reopt triggered"),
        (True,  "POWER_WASTE — ideal (training values)",
         "POWER_WASTE","POWER_WASTE","100%","Correct when given training-distribution values"),
        (False, "POWER_WASTE — simulator-realistic (1% load)",
         "POWER_WASTE","UNDERLOAD","100%","BUG: simulator power at low load=240W (<500W threshold)\nLinear power model cannot generate high-power + low-UE combo"),
    ]

    y0 = 2.28
    for ri, (ok, name, exp, pred, conf, note) in enumerate(results):
        yy = y0 + ri * 0.68
        bg = BG_NORMAL if ok else BG_SINR
        fc_res = GREEN if ok else RED
        rect(s, 0.3,  yy, 0.72, 0.62, bg,  "PASS" if ok else "FAIL",
             fs=9, bold=True, fc=fc_res)
        rect(s, 1.05, yy, 3.5, 0.62, LGRAY if ri%2==0 else WHITE,
             name, fs=9, fc=NAVY, align=PP_ALIGN.LEFT)
        rect(s, 4.6,  yy, 1.5, 0.62, LGRAY if ri%2==0 else WHITE,
             exp, fs=9, fc=NAVY)
        rect(s, 6.15, yy, 1.5, 0.62,
             BG_NORMAL if exp==pred else BG_SINR,
             pred, fs=9, bold=True, fc=GREEN if exp==pred else RED)
        rect(s, 7.7,  yy, 0.9, 0.62, LGRAY if ri%2==0 else WHITE,
             conf, fs=9, bold=True, fc=BLUE)
        rect(s, 8.65, yy, 4.4, 0.62, LGRAY if ri%2==0 else WHITE,
             note, fs=8, fc=DGRAY, align=PP_ALIGN.LEFT)

    for x, w, h in [(0.3,0.72,""), (1.05,3.5,"Scenario"),
                    (4.6,1.5,"Expected"), (6.15,1.5,"LSTM Predicted"),
                    (7.7,0.9,"Conf"), (8.65,4.4,"Finding")]:
        if h:
            rect(s, x, 2.08, w, 0.22, BLUE, h, fs=8, bold=True, fc=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 8 — Real Malleswaram Cell Health Snapshot
# ─────────────────────────────────────────────────────────────────────────────
def s8_snapshot(prs):
    s = blank(prs)
    hdr(s, "Real Deployment — 30-Cell Malleswaram Health Snapshot",
        "Post-KPI-agent rebalancing  |  DU-MLS-1 @ 88% load  |  DU-MLS-2 @ 72%  |  DU-MLS-3 @ 62%")
    pg(s, 8)

    # DU summary boxes
    dus = [
        ("DU-MLS-1", "12 cells", "88% load", ORANGE,  BG_OVER,
         "Heavy — was overloaded.\nKPI agent moved 4 cells away."),
        ("DU-MLS-2", "10 cells", "72% load", TEAL,    BG_UNDER,
         "Medium — received\nsome rebalanced cells."),
        ("DU-MLS-3", "8 cells",  "62% load", GREEN,   BG_NORMAL,
         "Light — received\nsome rebalanced cells."),
    ]
    for i, (du, cells, load, fg, bg, note) in enumerate(dus):
        x = 0.3 + i * 4.35
        rect(s, x, 1.05, 4.0, 0.5, fg, du, fs=13, bold=True)
        box(s, x, 1.6, 4.0, 1.5, bg, fc=fg, lc=fg)
        lbl(s, x+0.1, 1.65, 3.8, 0.4, cells + "  |  " + load, fs=11, bold=True, fc=fg)
        lbl(s, x+0.1, 2.1,  3.8, 0.9, note, fs=10, fc=DGRAY, align=PP_ALIGN.LEFT)

    # Selected cell details
    lbl(s, 0.3, 3.3, 12.7, 0.38,
        "Sample cell health at their DU load  (one simulated tick per cell):",
        fs=11, bold=True, fc=NAVY)

    cells_sample = [
        ("MLS_RWS_01", "DU-MLS-1", "5G n78", "88%", "8.2",  "792", "887W", "0.38%", "OVERLOAD", ORANGE, BG_OVER),
        ("MLS_18C_01", "DU-MLS-1", "5G n78", "83%", "10.1", "747", "851W", "0.22%", "NORMAL",   GREEN,  BG_NORMAL),
        ("MLS_SPG_01", "DU-MLS-2", "5G n78", "68%", "13.5", "612", "719W", "0.00%", "NORMAL",   GREEN,  BG_NORMAL),
        ("MLS_3MN_02", "DU-MLS-2", "4G B40", "67%", "15.2", "200", "133W", "0.00%", "NORMAL",   GREEN,  BG_NORMAL),
        ("MLS_6CR_01", "DU-MLS-3", "5G n78", "57%", "15.8", "513", "649W", "0.00%", "NORMAL",   GREEN,  BG_NORMAL),
        ("MLS_MGR_02", "DU-MLS-3", "4G B40", "60%", "13.7", "180", "128W", "0.00%", "NORMAL",   GREEN,  BG_NORMAL),
    ]

    hdrs = ["Cell ID", "DU", "Type", "PRB%", "SINR dB", "UEs", "Power", "Pkt Loss", "LSTM Class"]
    ws   = [1.9, 1.3, 1.0, 0.8, 0.9, 0.7, 0.85, 0.85, 1.5]
    y0 = 3.72
    xpos = 0.3
    for h, w in zip(hdrs, ws):
        rect(s, xpos, y0, w, 0.3, NAVY, h, fs=8, bold=True, fc=WHITE)
        xpos += w

    for ri, row in enumerate(cells_sample):
        yy = y0 + 0.3 + ri * 0.52
        bg = LGRAY if ri % 2 == 0 else WHITE
        *vals, cls, fg, cbg = row
        xpos = 0.3
        for vi, (v, w) in enumerate(zip(vals, ws[:-1])):
            rect(s, xpos, yy, w, 0.5, bg, v, fs=9,
                 bold=(vi == 0), fc=NAVY, align=PP_ALIGN.LEFT if vi==0 else PP_ALIGN.CENTER)
            xpos += w
        rect(s, xpos, yy, ws[-1], 0.5, cbg, cls, fs=9, bold=True, fc=fg)

    lbl(s, 0.3, 7.12, 12.7, 0.3,
        "At 88% load, some DU-MLS-1 cells breach 85% PRB threshold and trigger auto-move. "
        "DU-MLS-2 and DU-MLS-3 cells stay in NORMAL range after rebalancing.",
        fs=9, fc=DGRAY, italic=True)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 9 — How It Was Trained
# ─────────────────────────────────────────────────────────────────────────────
def s9_training(prs):
    s = blank(prs)
    hdr(s, "How the LSTM Was Trained — Synthetic Data, Realistic Distribution",
        "No live network data needed at training time. Synthetic sequences generated from 5G and 4G KPI profiles.")
    pg(s, 9)

    # Class distribution — big visual bars
    lbl(s, 0.3, 1.05, 6.5, 0.42, "Training Samples per Class", fs=13, bold=True, fc=NAVY)
    dist = [
        ("NORMAL",      3500, 70, GREEN,  BG_NORMAL),
        ("OVERLOAD",     750, 15, ORANGE, BG_OVER),
        ("UNDERLOAD",    400,  8, BLUE,   BG_UNDER),
        ("SINR LOW",     250,  5, RED,    BG_SINR),
        ("POWER WASTE",  100,  2, PURPLE, BG_PWR),
    ]
    for i, (name, n, pct, fg, bg) in enumerate(dist):
        y = 1.55 + i * 0.75
        rect(s, 0.3, y, 2.0, 0.62, bg, name, fs=11, bold=True, fc=fg)
        bw = max(pct / 70 * 5.8, 0.5)
        rect(s, 2.35, y, bw, 0.62, fg, f"{n:,} samples  ({pct}%)",
             fs=10, bold=True, fc=WHITE)

    lbl(s, 0.3, 5.42, 8.3, 0.38,
        "Rare classes (POWER_WASTE = 2%) use Weighted Sampler so the model learns them too",
        fs=9, fc=DGRAY, italic=True)

    # Right panel — training facts
    box(s, 8.8, 1.05, 4.2, 4.85, LGRAY, lc=BLUE)
    facts = [
        ("Epochs",         "60"),
        ("Batch size",     "256"),
        ("Optimiser",      "Adam"),
        ("LR schedule",    "Cosine annealing"),
        ("5G profile",     "n78, 64T64R, 950W"),
        ("4G profile",     "B3/B40, 4T4R, 200W"),
        ("Sequence drift", "Slow drift between steps"),
        ("Train / Val",    "80% / 20% split"),
        ("Saved to",       "kpi_model.pt"),
        ("Train time",     "~2 min on CPU"),
    ]
    lbl(s, 8.8, 1.08, 4.2, 0.42, "Training Configuration", fs=12, bold=True, fc=BLUE)
    for i, (k, v) in enumerate(facts):
        yy = 1.58 + i * 0.43
        bg = LGRAY if i % 2 == 0 else WHITE
        rect(s, 8.8,  yy, 2.1, 0.4, bg, k, fs=9, fc=NAVY, align=PP_ALIGN.LEFT)
        rect(s, 10.9, yy, 2.1, 0.4, bg, v, fs=9, fc=BLUE, bold=True)

    # What each sequence looks like
    box(s, 0.3, 5.88, 12.7, 1.2, WHITE, lc=TEAL)
    lbl(s, 0.5, 5.96, 12.3, 0.42,
        "Each training sequence = 6 timesteps with realistic temporal drift",
        fs=11, bold=True, fc=TEAL)
    lbl(s, 0.5, 6.42, 12.3, 0.55,
        "OVERLOAD example:  PRB [91%, 93%, 94%, 93%, 95%, 94%]  "
        "SINR [11.2, 10.8, 11.1, 10.5, 9.8, 10.2]  UEs [705, 718, 712, 722, 715, 720]\n"
        "Small drift between steps teaches the LSTM to detect trends, not just memorise snapshots.",
        fs=9, fc=DGRAY, align=PP_ALIGN.LEFT)


# ─────────────────────────────────────────────────────────────────────────────
# Slide 10 — 3 Known Gaps + Summary
# ─────────────────────────────────────────────────────────────────────────────
def s10_gaps(prs):
    s = blank(prs)
    hdr(s, "3 Known Gaps + Summary",
        "The LSTM works well for 5G cells. These gaps were found during validation testing.")
    pg(s, 10)

    gaps = [
        (RED,  "GAP 1 — POWER_WASTE unreachable in simulator",
         "The simulator uses a linear power model: power scales with load. "
         "At low load (1% = 9 UEs), power = 240W. "
         "The rule requires power > 500W to fire POWER_WASTE. "
         "The LSTM was trained on 880W + 8 UEs — it never sees this in real operation.",
         "Fix: change simulator to hold near-peak RF power when UEs < threshold "
         "(real mMIMO hardware does this)"),
        (ORANGE, "GAP 2 — 4G cells misclassified",
         "Training data was mixed 5G/4G but OVERLOAD class assumed 5G-scale values "
         "(UEs=720, power=940W, throughput=3100 Mbps). "
         "A 4G cell at overload has UEs=230, power=195W, throughput=140 Mbps. "
         "LSTM sees this as SINR_LOW or UNDERLOAD — not OVERLOAD.",
         "Fix: separate training and inference paths for 4G vs 5G cells"),
        (BLUE, "GAP 3 — SINR_LOW needs real interference injection",
         "The simulator degrades SINR only through load. "
         "Co-channel interference (the main real-world cause of SINR_LOW) is not modelled. "
         "SINR_LOW at moderate load (the training distribution) cannot be generated naturally.",
         "Fix: add interference_db env variable per cell in the simulator"),
    ]

    for i, (clr, title, desc, fix) in enumerate(gaps):
        y = 1.05 + i * 1.92
        rect(s, 0.3, y, 12.7, 0.45, clr, title, fs=11, bold=True)
        rect(s, 0.3, y+0.45, 8.4, 1.42, LGRAY, desc, fs=9, fc=DGRAY,
             align=PP_ALIGN.LEFT)
        box(s, 8.75, y+0.45, 4.25, 1.42, WHITE,
            "Recommended fix:\n" + fix, fs=9, fc=clr,
            align=PP_ALIGN.LEFT, lc=clr)

    # Summary strip
    rect(s, 0.3, 6.82, 12.7, 0.58, NAVY,
         "LSTM correctly classifies 5G NORMAL, OVERLOAD, and SINR_LOW (interference) at 100% confidence.  "
         "4G cells and POWER_WASTE need training data improvements.",
         fs=10, bold=True, fc=WHITE)


# ─────────────────────────────────────────────────────────────────────────────
def main():
    prs = new_prs()
    s1_title(prs)
    s2_analogy(prs)
    s3_kpis(prs)
    s4_memory(prs)
    s5_bidir(prs)
    s6_decisions(prs)
    s7_results(prs)
    s8_snapshot(prs)
    s9_training(prs)
    s10_gaps(prs)
    prs.save(OUT)
    print(f"Saved -> {OUT}")
    print("10 slides: Title, Analogy, 9 KPIs, Memory/Gates, Bidirectional, "
          "5 Decisions, Validation Results, 30-Cell Snapshot, Training, Gaps+Summary")


if __name__ == "__main__":
    main()
