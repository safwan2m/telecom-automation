#!/usr/bin/env python3
"""Generate C2 Query Flow PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree

NAVY    = RGBColor(0x0A, 0x29, 0x55)
BLUE    = RGBColor(0x17, 0x5C, 0xA6)
TEAL    = RGBColor(0x00, 0x7E, 0x8A)
GREEN   = RGBColor(0x1B, 0x6F, 0x42)
ORANGE  = RGBColor(0xC4, 0x50, 0x08)
PURPLE  = RGBColor(0x5A, 0x23, 0x8C)
DGRAY   = RGBColor(0x44, 0x44, 0x44)
LGRAY   = RGBColor(0xF0, 0xF2, 0xF5)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
L_ORG   = RGBColor(0xFF, 0xF0, 0xE8)
L_PUR   = RGBColor(0xF2, 0xEB, 0xFF)
L_TEAL  = RGBColor(0xE0, 0xF5, 0xF7)
QARROW  = BLUE
RARROW  = GREEN


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(slide, x, y, w, h, fill, border=None, bw=0.5):
    shp = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border:
        shp.line.color.rgb = border
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp

def box(slide, x, y, w, h, fill, lines, fsize=11, fc=WHITE, bold0=True):
    shp = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shp.line.width = Pt(0.5)
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(fsize)
        r.font.color.rgb = fc
        r.font.bold = (bold0 and i == 0)
    return shp

def lbl(slide, x, y, w, h, text, fsize=9, fc=DGRAY, bold=False, italic=False,
        align=PP_ALIGN.CENTER):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fsize)
    r.font.color.rgb = fc
    r.font.bold = bold
    r.font.italic = italic
    return txb

def arrow(slide, x1, y1, x2, y2, color=None, width=1.75, dashed=False):
    if color is None:
        color = QARROW
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    try:
        ln = conn.line._ln
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('len', 'med')
    except Exception:
        pass
    return conn


prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
s1 = blank(prs)
rect(s1, 0, 0, 13.33, 7.5, NAVY)
rect(s1, 0, 0, 0.22, 7.5, RGBColor(0x00, 0xB4, 0xF0))
lbl(s1, 0.45, 2.4,  12.0, 1.3, 'C2 Query Flow',
    fsize=52, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)
lbl(s1, 0.45, 3.85, 12.0, 0.9,
    '"Show details of MLS_RWS_01."',
    fsize=26, fc=RGBColor(0x99, 0xCC, 0xFF), italic=True, align=PP_ALIGN.LEFT)
lbl(s1, 0.45, 5.0,  12.0, 0.5,
    'Tool: query_cell(cell_id="MLS_RWS_01")  →  full config + 30-min KPI time series',
    fsize=13, fc=RGBColor(0x55, 0xDD, 0x88), bold=True, align=PP_ALIGN.LEFT)
lbl(s1, 0.45, 5.7,  12.0, 0.4,
    '5G NR  ·  n78 (3500 MHz)  ·  Nokia AirScale MAA 64T64R  ·  DU-MLS-1  ·  CU-MLS',
    fsize=11, fc=RGBColor(0x77, 0x99, 0xBB), align=PP_ALIGN.LEFT)

# ── SLIDE 2: Flow Diagram ─────────────────────────────────────────────────────
s2 = blank(prs)
rect(s2, 0, 0, 13.33, 7.5, LGRAY)
rect(s2, 0, 0, 13.33, 0.9, NAVY)
lbl(s2, 0.3, 0.05, 9.5, 0.8, 'C2 Query Flow — System Data Path',
    fsize=22, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)
lbl(s2, 10.0, 0.1, 3.0, 0.7, 'query_cell',
    fsize=15, fc=RGBColor(0x88, 0xCC, 0xFF), bold=True)

rect(s2, 0.95, 1.0,  12.35, 1.8,  L_ORG,  border=RGBColor(0xDD, 0xAA, 0x80))
rect(s2, 0.95, 2.8,  12.35, 1.95, L_PUR,  border=RGBColor(0xAA, 0x88, 0xCC))
rect(s2, 0.95, 4.75, 12.35, 2.65, L_TEAL, border=RGBColor(0x77, 0xBB, 0xBB))
lbl(s2, 0.0, 1.35, 0.98, 0.9, 'USER', fsize=8, fc=ORANGE, bold=True)
lbl(s2, 0.0, 3.1,  0.98, 0.9, 'AI',   fsize=8, fc=PURPLE, bold=True)
lbl(s2, 0.0, 5.4,  0.98, 0.9, 'DATA', fsize=8, fc=TEAL,   bold=True)

box(s2, 1.1, 1.15, 1.5, 0.85, ORANGE, ['User', 'query input'])
box(s2, 3.3, 1.15, 1.9, 0.85, BLUE,   ['chat.py', 'localhost:8082'])
box(s2, 1.1, 2.95, 2.2, 0.85, NAVY,   ['Orchestrator', 'port 8082'])
box(s2, 4.3, 2.95, 2.6, 0.85, PURPLE, ['Gemini LLM', 'gemini-2.5-flash'])
box(s2, 1.1, 4.9,  2.2, 0.85, TEAL,   ['Controller', 'port 8080'])
box(s2, 4.1, 5.85, 2.3, 0.85, GREEN,
    ['topology.json', 'PCI, band, vendor,', 'power, max_ues, DU/CU'], fsize=9)
# InfluxDB — C2 returns full 6-KPI time series
box(s2, 7.3, 5.85, 2.7, 0.85, GREEN,
    ['InfluxDB', '30-min series: UEs, PRB,', 'SINR, throughput, HO, loss'], fsize=9)

arrow(s2, 2.6,  1.575, 3.3,  1.575, color=QARROW)
lbl(s2, 2.6, 1.2, 0.7, 0.35, 'HTTP\nPOST', fsize=7, fc=BLUE)
arrow(s2, 4.25, 2.0,   2.2,  2.95,  color=QARROW)
lbl(s2, 2.7, 2.35, 1.5, 0.35, 'POST /chat', fsize=7, fc=BLUE)
arrow(s2, 3.3,  3.1,   4.3,  3.1,   color=QARROW)
lbl(s2, 3.3, 2.75, 1.05, 0.35, 'prompt +\nsnapshot', fsize=7, fc=BLUE)
arrow(s2, 4.3,  3.55,  3.3,  3.55,  color=RARROW)
lbl(s2, 3.2, 3.6, 1.15, 0.35, 'fn call:\nquery_cell', fsize=7, fc=GREEN)
arrow(s2, 2.2,  3.8,   2.2,  4.9,   color=QARROW)
lbl(s2, 2.25, 4.2, 1.6, 0.5, 'GET /cells/\nMLS_RWS_01', fsize=7, fc=BLUE)
arrow(s2, 2.5,  5.75,  4.1,  6.1,   color=QARROW)
lbl(s2, 2.65, 5.5, 1.4, 0.3, 'read full config', fsize=7, fc=BLUE)
arrow(s2, 3.0,  5.6,   7.3,  5.95,  color=QARROW)
lbl(s2, 4.6, 5.33, 2.5, 0.3, 'Flux query: 30-min KPI series', fsize=7, fc=BLUE)
arrow(s2, 4.6,  5.85,  3.1,  5.0,   color=RARROW, dashed=True)
lbl(s2, 3.55, 5.28, 1.2, 0.3, 'full config', fsize=7, fc=GREEN)
arrow(s2, 7.5,  5.85,  3.2,  5.2,   color=RARROW, dashed=True)
lbl(s2, 5.0, 5.65, 2.5, 0.3, '6-KPI time series', fsize=7, fc=GREEN)
arrow(s2, 1.9,  4.9,   1.9,  3.8,   color=RARROW, dashed=True)
lbl(s2, 0.05, 4.2, 1.8, 0.5, 'config +\nKPI series', fsize=7, fc=GREEN)
arrow(s2, 3.3,  3.65,  2.3,  2.1,   color=RARROW, dashed=True)
arrow(s2, 2.25, 2.05,  1.85, 1.75,  color=RARROW, dashed=True)
lbl(s2, 0.6, 2.5, 1.4, 0.35, 'NL reply', fsize=7, fc=GREEN)

rect(s2, 10.0, 1.15, 3.1, 1.45, WHITE, border=RGBColor(0xCC, 0xCC, 0xCC))
lbl(s2, 10.05, 1.18, 3.0, 0.35, 'Legend', fsize=10, fc=NAVY, bold=True)
arrow(s2, 10.1, 1.72, 10.7, 1.72, color=QARROW, width=1.5)
lbl(s2, 10.75, 1.58, 2.2, 0.3, 'Query / Request path', fsize=9, fc=DGRAY)
arrow(s2, 10.1, 2.09, 10.7, 2.09, color=RARROW, width=1.5, dashed=True)
lbl(s2, 10.75, 1.96, 2.2, 0.3, 'Response / Return path', fsize=9, fc=DGRAY)

# Gemini note — C2: formats full config + KPI trends (richest response)
box(s2, 7.8, 2.95, 3.0, 0.85, RGBColor(0x44, 0x22, 0x70),
    ['Formats all config fields', '+ summarises KPI trends', '→ full cell profile'],
    fsize=9, bold0=False)
arrow(s2, 6.9, 3.375, 7.8, 3.375, color=PURPLE, width=1.5)
lbl(s2, 6.9, 3.1, 0.9, 0.3, 'tool\nresult', fsize=7, fc=PURPLE)

# ── SLIDE 3: Step-by-step ─────────────────────────────────────────────────────
s3 = blank(prs)
rect(s3, 0, 0, 13.33, 7.5, LGRAY)
rect(s3, 0, 0, 13.33, 0.9, NAVY)
lbl(s3, 0.3, 0.05, 10.0, 0.8, 'C2 Query — Step-by-Step Trace',
    fsize=22, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)

steps = [
    ('1',  'User',         'Types query into chat.py'),
    ('2',  'chat.py',      'HTTP POST {session, message} to Orchestrator /chat'),
    ('3',  'Orchestrator', 'Injects live network snapshot into system prompt; calls Gemini'),
    ('4',  'Gemini LLM',   'Identifies cell MLS_RWS_01; emits query_cell(cell_id="MLS_RWS_01")'),
    ('5',  'Orchestrator', 'Dispatches query_cell → GET /cells/MLS_RWS_01 on Controller'),
    ('6',  'Controller',   'Reads topology.json: PCI=1, band=n78, vendor=Nokia, 64T64R, 3800 Mbps, DU-MLS-1'),
    ('7',  'Controller',   'Runs Flux query: 30-min series for connected_ues, PRB, SINR, throughput, HO rate, packet loss'),
    ('8',  'Controller',   'Returns merged full config + 30-min KPI time series for MLS_RWS_01'),
    ('9',  'Orchestrator', 'Passes tool result (cell detail) back to Gemini'),
    ('10', 'Gemini LLM',   'Formats all config fields and summarises 30-min KPI trends'),
    ('11', 'Gemini LLM',   'Generates detailed natural language cell profile for MLS_RWS_01'),
    ('12', 'Orchestrator', 'Returns HTTP response {reply} to chat.py'),
    ('13', 'chat.py',      'Prints full cell profile to user'),
]
col_colors = {'User': ORANGE, 'chat.py': BLUE, 'Orchestrator': NAVY,
              'Gemini LLM': PURPLE, 'Controller': TEAL}
ROW_H, START_Y = 0.46, 1.0
for i, (step, actor, action) in enumerate(steps):
    y = START_Y + i * ROW_H
    bg = RGBColor(0xFF, 0xFF, 0xFF) if i % 2 == 0 else RGBColor(0xE8, 0xEC, 0xF2)
    rect(s3, 0.3, y, 12.7, ROW_H, bg, border=RGBColor(0xCC, 0xCC, 0xCC), bw=0.3)
    lbl(s3, 0.35, y+0.05, 0.55, ROW_H-0.1, step, fsize=11, fc=NAVY, bold=True)
    rect(s3, 0.95, y+0.08, 1.6, ROW_H-0.16, col_colors.get(actor, DGRAY),
         border=RGBColor(0xAA, 0xAA, 0xAA), bw=0.3)
    lbl(s3, 0.95, y+0.08, 1.6, ROW_H-0.16, actor, fsize=9, fc=WHITE, bold=True)
    lbl(s3, 2.65, y+0.05, 10.3, ROW_H-0.1, action,
        fsize=10, fc=DGRAY, align=PP_ALIGN.LEFT)

# ── SLIDE 4: MLS_RWS_01 Cell Profile ─────────────────────────────────────────
s4 = blank(prs)
rect(s4, 0, 0, 13.33, 7.5, LGRAY)
rect(s4, 0, 0, 13.33, 0.9, NAVY)
lbl(s4, 0.3, 0.05, 10.0, 0.8, 'C2 Response Profile — MLS_RWS_01',
    fsize=22, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)

# Config table (left)
rect(s4, 0.3, 1.0, 6.0, 0.5, TEAL)
lbl(s4, 0.3, 1.0, 6.0, 0.5, 'Static Config  (from topology.json)', fsize=11, fc=WHITE, bold=True)
config_rows = [
    ('Generation', '5G NR'),
    ('Band', 'n78  ·  3500 MHz'),
    ('PCI', '1'),
    ('Vendor', 'Nokia'),
    ('Hardware', 'AirScale MAA 64T64R'),
    ('Antenna', '64T64R'),
    ('Peak DL', '3800 Mbps'),
    ('TX Power', '1000 W'),
    ('Max UEs', '900'),
    ('DU', 'DU-MLS-1'),
    ('CU', 'CU-MLS'),
]
for i, (k, v) in enumerate(config_rows):
    y = 1.5 + i * 0.46
    bg = WHITE if i % 2 == 0 else RGBColor(0xE8, 0xEC, 0xF2)
    rect(s4, 0.3, y, 6.0, 0.46, bg, border=RGBColor(0xCC, 0xCC, 0xCC), bw=0.3)
    lbl(s4, 0.35, y+0.05, 2.2, 0.36, k, fsize=10, fc=NAVY, bold=True, align=PP_ALIGN.LEFT)
    lbl(s4, 2.6,  y+0.05, 3.6, 0.36, v, fsize=10, fc=DGRAY, align=PP_ALIGN.LEFT)

# KPI table (right)
rect(s4, 6.8, 1.0, 6.2, 0.5, GREEN)
lbl(s4, 6.8, 1.0, 6.2, 0.5, '30-min KPI Series  (from InfluxDB)', fsize=11, fc=WHITE, bold=True)
kpi_rows = [
    ('connected_ues',      'Active UEs on this cell over 30 min'),
    ('prb_dl_pct',         'Downlink PRB utilisation %'),
    ('sinr_db',            'Signal-to-Interference-plus-Noise Ratio'),
    ('dl_throughput_mbps', 'Downlink throughput (Mbps)'),
    ('ho_success_rate',    'Handover success rate'),
    ('packet_loss_pct',    'Packet loss percentage'),
]
for i, (k, v) in enumerate(kpi_rows):
    y = 1.5 + i * 0.55
    bg = WHITE if i % 2 == 0 else RGBColor(0xE8, 0xEC, 0xF2)
    rect(s4, 6.8, y, 6.2, 0.55, bg, border=RGBColor(0xCC, 0xCC, 0xCC), bw=0.3)
    lbl(s4, 6.85, y+0.06, 2.8, 0.43, k, fsize=10, fc=TEAL, bold=True, align=PP_ALIGN.LEFT)
    lbl(s4, 9.7,  y+0.06, 3.2, 0.43, v, fsize=10, fc=DGRAY, align=PP_ALIGN.LEFT)

out = r'c:\Users\gurs2\OneDrive\Documents\GitHub\telecom-automation\test\C2_Query_Flow.pptx'
prs.save(out)
print(f'Saved: {out}')
