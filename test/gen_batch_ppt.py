#!/usr/bin/env python3
"""Batch generator: F1, F2, A1, A2, O1, O2, I1, I2, E1, E2 PPT flow diagrams."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import os

OUT = r'c:\Users\gurs2\OneDrive\Documents\GitHub\telecom-automation\test'

# ─── Colors ──────────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0A, 0x29, 0x55)
BLUE   = RGBColor(0x17, 0x5C, 0xA6)
TEAL   = RGBColor(0x00, 0x7E, 0x8A)
GREEN  = RGBColor(0x1B, 0x6F, 0x42)
ORANGE = RGBColor(0xC4, 0x50, 0x08)
PURPLE = RGBColor(0x5A, 0x23, 0x8C)
RED    = RGBColor(0xB5, 0x1A, 0x1A)
DGRAY  = RGBColor(0x44, 0x44, 0x44)
LGRAY  = RGBColor(0xF0, 0xF2, 0xF5)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
L_ORG  = RGBColor(0xFF, 0xF0, 0xE8)
L_PUR  = RGBColor(0xF2, 0xEB, 0xFF)
L_TEAL = RGBColor(0xE0, 0xF5, 0xF7)
QA     = BLUE
RA     = GREEN
GOLD   = RGBColor(0xF5, 0xC5, 0x18)


# ─── Helpers ─────────────────────────────────────────────────────────────────

def new_prs():
    p = Presentation()
    p.slide_width  = Inches(13.33)
    p.slide_height = Inches(7.5)
    return p

def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def rect(s, x, y, w, h, fill, border=None, bw=0.5):
    shp = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if border: shp.line.color.rgb = border; shp.line.width = Pt(bw)
    else: shp.line.fill.background()
    return shp

def box(s, x, y, w, h, fill, lines, fs=11, fc=WHITE, b0=True):
    shp = s.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = RGBColor(0xCC,0xCC,0xCC); shp.line.width = Pt(0.5)
    tf = shp.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, t in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = t
        r.font.size = Pt(fs); r.font.color.rgb = fc; r.font.bold = (b0 and i == 0)
    return shp

def lbl(s, x, y, w, h, text, fs=9, fc=DGRAY, bold=False, italic=False, align=PP_ALIGN.CENTER):
    t = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = t.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(fs); r.font.color.rgb = fc; r.font.bold = bold; r.font.italic = italic
    return t

def arr(s, x1, y1, x2, y2, color=None, w=1.75, dashed=False):
    if color is None: color = QA
    c = s.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT,
                               Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    if dashed:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    try:
        ln = c.line._ln
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle'); tail.set('w', 'med'); tail.set('len', 'med')
    except Exception: pass
    return c


# ─── Shared slide builders ───────────────────────────────────────────────────

def title_slide(prs, qid, title_text, query_text, tool_text, caption=''):
    s = blank(prs)
    rect(s, 0, 0, 13.33, 7.5, NAVY)
    rect(s, 0, 0, 0.22, 7.5, RGBColor(0x00,0xB4,0xF0))
    lbl(s, 0.45, 2.4, 12.0, 1.3, f'{qid} Query Flow',
        fs=52, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)
    lbl(s, 0.45, 3.85, 12.0, 0.9, f'"{query_text}"',
        fs=18, fc=RGBColor(0x99,0xCC,0xFF), italic=True, align=PP_ALIGN.LEFT)
    lbl(s, 0.45, 5.0, 12.0, 0.5, f'Tool: {tool_text}',
        fs=13, fc=RGBColor(0x55,0xDD,0x88), bold=True, align=PP_ALIGN.LEFT)
    if caption:
        lbl(s, 0.45, 5.7, 12.0, 0.4, caption,
            fs=11, fc=RGBColor(0x77,0x99,0xBB), align=PP_ALIGN.LEFT)


def _layer_bands(s):
    rect(s, 0.95, 1.0,  12.35, 1.8,  L_ORG,  border=RGBColor(0xDD,0xAA,0x80))
    rect(s, 0.95, 2.8,  12.35, 1.95, L_PUR,  border=RGBColor(0xAA,0x88,0xCC))
    rect(s, 0.95, 4.75, 12.35, 2.65, L_TEAL, border=RGBColor(0x77,0xBB,0xBB))
    lbl(s, 0.0, 1.35, 0.98, 0.9, 'USER', fs=8, fc=ORANGE, bold=True)
    lbl(s, 0.0, 3.1,  0.98, 0.9, 'AI',   fs=8, fc=PURPLE, bold=True)
    lbl(s, 0.0, 5.4,  0.98, 0.9, 'DATA', fs=8, fc=TEAL,   bold=True)

def _legend(s):
    rect(s, 10.0, 1.15, 3.1, 1.45, WHITE, border=RGBColor(0xCC,0xCC,0xCC))
    lbl(s, 10.05, 1.18, 3.0, 0.35, 'Legend', fs=10, fc=NAVY, bold=True)
    arr(s, 10.1, 1.72, 10.7, 1.72, color=QA, w=1.5)
    lbl(s, 10.75, 1.58, 2.2, 0.3, 'Query / Request path', fs=9, fc=DGRAY)
    arr(s, 10.1, 2.09, 10.7, 2.09, color=RA, w=1.5, dashed=True)
    lbl(s, 10.75, 1.96, 2.2, 0.3, 'Response / Return path', fs=9, fc=DGRAY)

def _common_boxes(s):
    box(s, 1.1, 1.15, 1.5, 0.85, ORANGE, ['User', 'query input'])
    box(s, 3.3, 1.15, 1.9, 0.85, BLUE,   ['chat.py', 'localhost:8082'])
    box(s, 1.1, 2.95, 2.2, 0.85, NAVY,   ['Orchestrator', 'port 8082'])
    box(s, 4.3, 2.95, 2.6, 0.85, PURPLE, ['Gemini LLM', 'gemini-2.5-flash'])

def _common_user_arrows(s):
    arr(s, 2.6,  1.575, 3.3,  1.575)
    lbl(s, 2.6,  1.2,   0.7,  0.35, 'HTTP\nPOST', fs=7, fc=BLUE)
    arr(s, 4.25, 2.0,   2.2,  2.95)
    lbl(s, 2.7,  2.35,  1.5,  0.35, 'POST /chat', fs=7, fc=BLUE)
    arr(s, 3.3,  3.1,   4.3,  3.1)
    lbl(s, 3.3,  2.75,  1.05, 0.35, 'SYSTEM_PROMPT\n+ snapshot\n+ history', fs=7, fc=BLUE)

def _common_return_arrows(s):
    arr(s, 1.9, 4.9, 1.9, 3.8, color=RA, dashed=True)
    arr(s, 3.3, 3.65, 2.3, 2.1, color=RA, dashed=True)
    arr(s, 2.25, 2.05, 1.85, 1.75, color=RA, dashed=True)
    lbl(s, 0.6, 2.5, 1.4, 0.35, 'NL reply', fs=7, fc=GREEN)

def _data_layer(s, topo_lines=None, influx_lines=None):
    box(s, 1.1, 4.9, 2.2, 0.85, TEAL, ['Controller', 'port 8080'])
    tl = topo_lines or ['topology.json', 'cell/DU/CU config']
    il = influx_lines or ['InfluxDB', 'port 8086']
    box(s, 4.1, 5.85, 2.3, 0.85, GREEN, tl, fs=10)
    box(s, 7.3, 5.85, 2.3, 0.85, GREEN, il, fs=10)

def _standard_flow_slide(prs, header, tool_badge, fn_lbl, endpoint_lbl,
                          topo_lbl, influx_lbl, return_data_lbl, influx_return_lbl,
                          gemini_note, topo_lines=None, influx_lines=None):
    s = blank(prs)
    rect(s, 0, 0, 13.33, 7.5, LGRAY)
    rect(s, 0, 0, 13.33, 0.9, NAVY)
    lbl(s, 0.3, 0.05, 9.5, 0.8, header, fs=20, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)
    lbl(s, 10.0, 0.1, 3.0, 0.7, tool_badge, fs=14, fc=RGBColor(0x88,0xCC,0xFF), bold=True)
    _layer_bands(s)
    _common_boxes(s)
    _data_layer(s, topo_lines, influx_lines)
    _common_user_arrows(s)
    arr(s, 4.3, 3.55, 3.3, 3.55, color=RA)
    lbl(s, 3.25, 3.6, 1.1, 0.35, fn_lbl, fs=7, fc=GREEN)
    arr(s, 2.2, 3.8, 2.2, 4.9)
    lbl(s, 2.25, 4.25, 1.4, 0.45, endpoint_lbl, fs=7, fc=BLUE)
    arr(s, 2.5, 5.75, 4.1, 6.1)
    lbl(s, 2.7, 5.5, 1.3, 0.35, topo_lbl, fs=7, fc=BLUE)
    arr(s, 3.0, 5.6, 7.3, 5.95)
    lbl(s, 4.6, 5.33, 2.3, 0.3, influx_lbl, fs=7, fc=BLUE)
    arr(s, 4.6, 5.85, 3.1, 5.0, color=RA, dashed=True)
    lbl(s, 3.5, 5.28, 1.3, 0.3, return_data_lbl, fs=7, fc=GREEN)
    arr(s, 7.5, 5.85, 3.2, 5.2, color=RA, dashed=True)
    lbl(s, 5.0, 5.65, 2.3, 0.3, influx_return_lbl, fs=7, fc=GREEN)
    _common_return_arrows(s)
    lbl(s, 0.1, 4.2, 1.75, 0.45, 'merged\nresult', fs=7, fc=GREEN)
    _legend(s)
    box(s, 7.8, 2.95, 3.0, 0.85, RGBColor(0x44,0x22,0x70), gemini_note, fs=9, b0=False)
    arr(s, 6.9, 3.375, 7.8, 3.375, color=PURPLE, w=1.5)
    lbl(s, 6.9, 3.1, 0.9, 0.3, 'tool\nresult', fs=7, fc=PURPLE)


def no_tool_slide(prs, header, gemini_note_lines):
    """Simplified flow: no Controller or InfluxDB — pure LLM reasoning."""
    s = blank(prs)
    rect(s, 0, 0, 13.33, 7.5, LGRAY)
    rect(s, 0, 0, 13.33, 0.9, NAVY)
    lbl(s, 0.3, 0.05, 9.5, 0.8, header, fs=20, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)
    lbl(s, 10.0, 0.1, 3.1, 0.7, 'No tool call', fs=13, fc=GOLD, bold=True)
    # Only User + AI bands
    rect(s, 0.95, 1.0,  12.35, 1.8,  L_ORG, border=RGBColor(0xDD,0xAA,0x80))
    rect(s, 0.95, 2.8,  12.35, 3.5,  L_PUR, border=RGBColor(0xAA,0x88,0xCC))
    lbl(s, 0.0, 1.35, 0.98, 0.9, 'USER', fs=8, fc=ORANGE, bold=True)
    lbl(s, 0.0, 3.5,  0.98, 0.9, 'AI',   fs=8, fc=PURPLE, bold=True)
    box(s, 1.1, 1.15, 1.5, 0.85, ORANGE, ['User', 'query input'])
    box(s, 3.3, 1.15, 1.9, 0.85, BLUE,   ['chat.py', 'localhost:8082'])
    box(s, 1.1, 2.95, 2.2, 0.85, NAVY,   ['Orchestrator', 'port 8082'])
    box(s, 4.3, 2.95, 2.6, 0.85, PURPLE, ['Gemini LLM', 'gemini-2.5-flash'])
    # Conversation history box
    box(s, 8.0, 2.95, 4.9, 0.85, RGBColor(0x1A, 0x5A, 0x2A),
        ['Conversation History', '(plan results, tool call records,', 'prior alerts & moves)'], fs=9)
    arr(s, 6.9, 3.375, 8.0, 3.375, color=GREEN, w=1.5)
    lbl(s, 6.9, 3.1, 1.1, 0.3, 'reads\ncontext', fs=7, fc=GREEN)
    # Gemini reasoning box
    box(s, 4.3, 4.3, 2.6, 1.0, RGBColor(0x44,0x22,0x70), gemini_note_lines, fs=9, b0=False)
    arr(s, 5.6, 3.8, 5.6, 4.3, color=PURPLE, w=1.5)
    lbl(s, 5.65, 3.9, 1.2, 0.35, 'reasons\nfrom history', fs=7, fc=PURPLE)
    # Arrows
    arr(s, 2.6, 1.575, 3.3, 1.575)
    lbl(s, 2.6, 1.2, 0.7, 0.35, 'HTTP\nPOST', fs=7, fc=BLUE)
    arr(s, 4.25, 2.0, 2.2, 2.95)
    lbl(s, 2.7, 2.35, 1.5, 0.35, 'POST /chat', fs=7, fc=BLUE)
    arr(s, 3.3, 3.1, 4.3, 3.1)
    lbl(s, 3.3, 2.75, 1.05, 0.35, 'SYSTEM_PROMPT\n+ snapshot\n+ history', fs=7, fc=BLUE)
    arr(s, 4.3, 5.3, 3.3, 2.2, color=RA, dashed=True)
    arr(s, 2.25, 2.05, 1.85, 1.75, color=RA, dashed=True)
    lbl(s, 0.6, 2.5, 1.4, 0.35, 'NL reply', fs=7, fc=GREEN)


def multitool_slide(prs, header, tool1_lbl, ep1, tool2_lbl, ep2, gemini_note):
    """Flow for multi-tool queries (get_alerts/query_network + move_cell)."""
    s = blank(prs)
    rect(s, 0, 0, 13.33, 7.5, LGRAY)
    rect(s, 0, 0, 13.33, 0.9, NAVY)
    lbl(s, 0.3, 0.05, 9.5, 0.8, header, fs=20, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)
    lbl(s, 10.0, 0.1, 3.1, 0.7, 'Multi-tool', fs=13, fc=GOLD, bold=True)
    _layer_bands(s)
    _common_boxes(s)
    # Data layer: Controller + InfluxDB + topology.json
    box(s, 1.1, 4.9, 2.2, 0.75, TEAL,  ['Controller', 'port 8080'])
    box(s, 4.0, 5.8, 2.2, 0.75, GREEN, ['topology.json', 'cell assignments'], fs=10)
    box(s, 7.1, 5.8, 2.2, 0.75, GREEN, ['InfluxDB', 'port 8086'], fs=10)
    _common_user_arrows(s)
    # Tool call 1
    arr(s, 4.3, 3.2, 3.3, 3.2, color=RA)
    lbl(s, 3.2, 3.25, 1.1, 0.3, tool1_lbl, fs=7, fc=GREEN)
    arr(s, 2.2, 3.8, 2.2, 4.9)
    lbl(s, 2.25, 4.2, 1.4, 0.45, ep1, fs=7, fc=BLUE)
    arr(s, 2.5, 5.65, 4.0, 5.95)
    arr(s, 3.0, 5.5, 7.1, 5.8)
    arr(s, 4.3, 5.8, 3.1, 5.0, color=RA, dashed=True)
    arr(s, 7.3, 5.8, 3.2, 5.1, color=RA, dashed=True)
    lbl(s, 0.1, 4.1, 1.8, 0.5, 'tool 1\nresult', fs=7, fc=GREEN)
    arr(s, 1.9, 4.9, 1.9, 3.85, color=RA, dashed=True)
    # Tool call 2
    arr(s, 4.3, 3.55, 3.3, 3.55, color=RGBColor(0xDD,0x88,0x00))
    lbl(s, 3.2, 3.6, 1.1, 0.3, tool2_lbl, fs=7, fc=RGBColor(0xCC,0x77,0x00))
    arr(s, 2.0, 3.8, 2.0, 4.9, color=RGBColor(0xDD,0x88,0x00))
    lbl(s, 2.25, 4.5, 1.4, 0.35, ep2, fs=7, fc=RGBColor(0xCC,0x77,0x00))
    arr(s, 2.5, 5.7, 4.0, 6.0, color=RGBColor(0xDD,0x88,0x00))
    arr(s, 4.3, 6.0, 3.1, 5.05, color=RA, dashed=True)
    lbl(s, 0.1, 4.5, 1.8, 0.4, 'tool 2\nresult', fs=7, fc=GREEN)
    # Return to user
    arr(s, 3.3, 3.65, 2.3, 2.1, color=RA, dashed=True)
    arr(s, 2.25, 2.05, 1.85, 1.75, color=RA, dashed=True)
    lbl(s, 0.6, 2.5, 1.4, 0.35, 'NL reply', fs=7, fc=GREEN)
    _legend(s)
    box(s, 7.8, 2.95, 3.0, 0.85, RGBColor(0x44,0x22,0x70), gemini_note, fs=9, b0=False)
    arr(s, 6.9, 3.375, 7.8, 3.375, color=PURPLE, w=1.5)
    lbl(s, 6.9, 3.1, 0.9, 0.3, 'tool\nresults', fs=7, fc=PURPLE)
    # Multi-tool note
    rect(s, 10.0, 2.8, 3.1, 0.6, RGBColor(0xF5,0xC5,0x18), border=DGRAY, bw=0.5)
    lbl(s, 10.05, 2.82, 3.0, 0.56,
        'Orange arrows = Tool 2 (action)\nGreen arrows = Tool 1 (read)',
        fs=8, fc=DGRAY)


def steps_slide(prs, header, steps):
    s = blank(prs)
    rect(s, 0, 0, 13.33, 7.5, LGRAY)
    rect(s, 0, 0, 13.33, 0.9, NAVY)
    lbl(s, 0.3, 0.05, 10.0, 0.8, header, fs=22, fc=WHITE, bold=True, align=PP_ALIGN.LEFT)
    COL = {'User': ORANGE, 'chat.py': BLUE, 'Orchestrator': NAVY,
            'Gemini LLM': PURPLE, 'Controller': TEAL}
    RH, SY = 0.44, 1.0
    for i, (step, actor, action) in enumerate(steps):
        y = SY + i * RH
        bg = WHITE if i % 2 == 0 else RGBColor(0xE8,0xEC,0xF2)
        rect(s, 0.3, y, 12.7, RH, bg, border=RGBColor(0xCC,0xCC,0xCC), bw=0.3)
        lbl(s, 0.35, y+0.04, 0.55, RH-0.08, step, fs=11, fc=NAVY, bold=True)
        rect(s, 0.95, y+0.07, 1.6, RH-0.14, COL.get(actor, DGRAY),
             border=RGBColor(0xAA,0xAA,0xAA), bw=0.3)
        lbl(s, 0.95, y+0.07, 1.6, RH-0.14, actor, fs=9, fc=WHITE, bold=True)
        lbl(s, 2.65, y+0.04, 10.3, RH-0.08, action,
            fs=10, fc=DGRAY, align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════════════
# F1 — No tool, pure LLM reasoning from plan context
# ═══════════════════════════════════════════════════════════════════════════════
def make_f1():
    prs = new_prs()
    title_slide(prs, 'F1',
        'Why was that site selected for the new cell?',
        'No tool call  →  LLM explains from plan context',
        'Context: follow-up after plan_network')
    no_tool_slide(prs, 'F1 Query Flow — Plan Explanation',
        ['Reads plan from context:', 'density_weight + budget', '→ explains site selection'])
    steps_slide(prs, 'F1 Query — Step-by-Step Trace', [
        ('1',  'User',         'Types follow-up query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session_id, message} → Orchestrator /chat  (pure stdlib urllib; --session flag)'),
        ('3',  'Orchestrator', 'Calls build_network_context() → GET /network; appends snapshot to SYSTEM_PROMPT'),
        ('4',  'Orchestrator', 'Sends SYSTEM_PROMPT + snapshot + full session history (types.Content) to Gemini'),
        ('5',  'Gemini LLM',   'No tool call — reads plan_network result from session history'),
        ('6',  'Gemini LLM',   'Identifies candidate site density_weight and budget fit'),
        ('7',  'Gemini LLM',   'Explains SINR margin, PCI assignment, DU selection logic; loop exits (no tool calls)'),
        ('8',  'Orchestrator', 'Streams text/plain chunks via sync generator (Starlette thread pool → StreamingResponse)'),
        ('9',  'chat.py',      'Prints streamed explanation to operator terminal'),
    ])
    prs.save(os.path.join(OUT, 'F1_Query_Flow.pptx'))
    print('Saved F1_Query_Flow.pptx')


# ═══════════════════════════════════════════════════════════════════════════════
# F2 — query_network to cross-reference plan cells with current topology
# ═══════════════════════════════════════════════════════════════════════════════
def make_f2():
    prs = new_prs()
    title_slide(prs, 'F2',
        'What cells are affected by the new deployment?',
        'query_network  →  cross-reference plan cells vs current topology',
        'Context: follow-up after plan_network')
    _standard_flow_slide(prs,
        'F2 Query Flow — Affected Cells Analysis', 'query_network',
        'fn call:\nquery_network', 'GET /network',
        'read config', 'Flux query: cell KPIs',
        'cell config', 'KPIs per cell',
        ['Cross-ref plan cells', 'vs current topology', '→ list affected cells'])
    steps_slide(prs, 'F2 Query — Step-by-Step Trace', [
        ('1',  'User',         'Types follow-up query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session_id, message} → Orchestrator /chat  (pure stdlib urllib; --session flag)'),
        ('3',  'Orchestrator', 'Calls build_network_context() → GET /network; appends snapshot to SYSTEM_PROMPT; sends full session history to Gemini'),
        ('4',  'Gemini LLM',   'Emits query_network() to get current topology for cross-referencing'),
        ('5',  'Orchestrator', 'Dispatches query_network → GET /network on Controller; yields *[calling tool...]* inline'),
        ('6',  'Controller',   'Reads topology.json + InfluxDB KPIs; returns 30-cell snapshot'),
        ('7',  'Orchestrator', 'JSON-sanitises result; appends FunctionResponse to session history; re-calls Gemini (while True loop)'),
        ('8',  'Gemini LLM',   'Reads new plan cells from session history'),
        ('9',  'Gemini LLM',   'Cross-references plan cells vs existing cells (proximity, PCI); loop exits'),
        ('10', 'Gemini LLM',   'Generates list of affected cells with reasoning'),
        ('11', 'Orchestrator', 'Streams text/plain chunks via sync generator (StreamingResponse)'),
        ('12', 'chat.py',      'Prints streamed affected cell list to operator terminal'),
    ])
    prs.save(os.path.join(OUT, 'F2_Query_Flow.pptx'))
    print('Saved F2_Query_Flow.pptx')


# ═══════════════════════════════════════════════════════════════════════════════
# A1 — get_alerts, last 60 min, all severities
# ═══════════════════════════════════════════════════════════════════════════════
def make_a1():
    prs = new_prs()
    title_slide(prs, 'A1',
        'What anomalies currently exist in the network?',
        'get_alerts(minutes=60)  →  all severities, all types',
        '4G/5G NSA  ·  KPI Agent LSTM classifier  ·  InfluxDB alerts measurement')
    _standard_flow_slide(prs,
        'A1 Query Flow — Network Anomaly Detection', 'get_alerts',
        'fn call:\nget_alerts', 'GET /alerts\n?minutes=60',
        'N/A (alerts only)', 'Flux: alerts last 60 min',
        'N/A', 'alert records',
        ['Groups by severity', '& alert_type', '→ anomaly summary'],
        topo_lines=['topology.json', '(context only)'],
        influx_lines=['InfluxDB alerts', 'OVERLOAD/SINR_LOW/\nPOWER_WASTE/UNDERLOAD'])
    steps_slide(prs, 'A1 Query — Step-by-Step Trace', [
        ('1',  'User',         'Types query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session_id, message} → Orchestrator /chat  (pure stdlib urllib; --session flag)'),
        ('3',  'Orchestrator', 'Calls build_network_context() → GET /network; appends snapshot to SYSTEM_PROMPT; sends to Gemini'),
        ('4',  'Gemini LLM',   'Emits get_alerts(minutes=60) — all severities and types'),
        ('5',  'Orchestrator', 'Dispatches get_alerts → queries InfluxDB directly via Flux; yields *[calling tool...]* inline'),
        ('6',  'Controller',   'Flux query on InfluxDB alerts measurement: last 60 min, all severities'),
        ('7',  'InfluxDB',     'Returns alert records written by KPI Agent LSTM classifier'),
        ('8',  'Orchestrator', 'JSON-sanitises result; appends FunctionResponse to session history; re-calls Gemini (while True loop)'),
        ('9',  'Gemini LLM',   'Groups alerts by severity (CRITICAL/WARNING/INFO) and type; loop exits'),
        ('10', 'Gemini LLM',   'Generates natural language anomaly summary'),
        ('11', 'Orchestrator', 'Streams text/plain chunks via sync generator (StreamingResponse)'),
        ('12', 'chat.py',      'Prints streamed anomaly report to operator terminal'),
    ])
    prs.save(os.path.join(OUT, 'A1_Query_Flow.pptx'))
    print('Saved A1_Query_Flow.pptx')


# ═══════════════════════════════════════════════════════════════════════════════
# A2 — get_alerts(OVERLOAD) or query_network PRB filter
# ═══════════════════════════════════════════════════════════════════════════════
def make_a2():
    prs = new_prs()
    title_slide(prs, 'A2',
        'Identify overloaded cells.',
        'get_alerts(alert_type=OVERLOAD)  OR  query_network → filter PRB > 80%',
        '4G/5G NSA  ·  Two-path: KPI Agent alert vs direct PRB check')
    _standard_flow_slide(prs,
        'A2 Query Flow — Overload Detection (Path A: get_alerts)', 'get_alerts',
        'fn call:\nget_alerts\n(OVERLOAD)', 'GET /alerts\n?type=OVERLOAD',
        'N/A', 'Flux: OVERLOAD alerts',
        'N/A', 'OVERLOAD alert records',
        ['Identifies cells', 'with OVERLOAD alert', '→ lists with confidence'],
        topo_lines=['topology.json', '(context only)'],
        influx_lines=['InfluxDB alerts', 'filter: OVERLOAD\n+ confidence ≥ 0.70'])
    steps_slide(prs, 'A2 Query — Step-by-Step Trace (Path A)', [
        ('1',  'User',         'Types query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session_id, message} → Orchestrator /chat  (pure stdlib urllib; --session flag)'),
        ('3',  'Orchestrator', 'Calls build_network_context() → GET /network; appends snapshot to SYSTEM_PROMPT; sends to Gemini'),
        ('4',  'Gemini LLM',   'Emits get_alerts(alert_type="OVERLOAD", minutes=60)'),
        ('5',  'Orchestrator', 'Dispatches get_alerts → direct Flux query on InfluxDB; yields *[calling tool...]* inline'),
        ('6',  'InfluxDB',     'Returns OVERLOAD records with cell_id, severity, confidence ≥ 0.70'),
        ('7',  'Orchestrator', 'JSON-sanitises result; appends FunctionResponse to session history; re-calls Gemini (while True loop)'),
        ('8',  'Gemini LLM',   'Lists overloaded cells with severity and confidence score; loop exits'),
        ('9',  'Gemini LLM',   'Generates natural language overload report'),
        ('10', 'Orchestrator', 'Streams text/plain chunks via sync generator (StreamingResponse)'),
        ('11', 'chat.py',      'Prints streamed overloaded cell list to operator terminal'),
    ])
    prs.save(os.path.join(OUT, 'A2_Query_Flow.pptx'))
    print('Saved A2_Query_Flow.pptx')


# ═══════════════════════════════════════════════════════════════════════════════
# O1 — get_alerts + move_cell (multi-tool)
# ═══════════════════════════════════════════════════════════════════════════════
def make_o1():
    prs = new_prs()
    title_slide(prs, 'O1',
        'Optimize the network for load balancing.',
        'get_alerts(OVERLOAD)  →  move_cell × N  (multi-tool)',
        '4G/5G NSA  ·  SON Agent  ·  Autonomous load rebalancing')
    multitool_slide(prs,
        'O1 Query Flow — Load Balancing (Multi-Tool)',
        'fn call 1:\nget_alerts\n(OVERLOAD)', 'GET /alerts\n(OVERLOAD)',
        'fn call 2:\nmove_cell\n(×N cells)', 'POST /move/cell',
        ['get_alerts → overloaded', 'cells; move_cell each', 'to lightest DU'])
    steps_slide(prs, 'O1 Query — Step-by-Step Trace', [
        ('1',  'User',         'Types query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session_id, message} → Orchestrator /chat  (pure stdlib urllib; --session flag)'),
        ('3',  'Orchestrator', 'Calls build_network_context() → GET /network; appends snapshot to SYSTEM_PROMPT; sends to Gemini'),
        ('4',  'Gemini LLM',   'Tool call 1: get_alerts(alert_type="OVERLOAD")  [while True loop, iteration 1]'),
        ('5',  'Orchestrator', 'Direct Flux query on InfluxDB: OVERLOAD alerts; yields *[calling tool: get_alerts...]* inline'),
        ('6',  'InfluxDB',     'Returns OVERLOAD records with cell_id, DU, severity, confidence'),
        ('7',  'Orchestrator', 'JSON-sanitises; appends FunctionResponse to session history; re-calls Gemini [loop iteration 2]'),
        ('8',  'Gemini LLM',   'Identifies lightest DU for each overloaded cell from network snapshot'),
        ('9',  'Gemini LLM',   'Tool call 2: move_cell(cell_id=X, target_du_id=Y)  [can batch multiple tool calls per turn]'),
        ('10', 'Orchestrator', 'POST /move/cell on Controller; yields *[calling tool: move_cell...]* inline'),
        ('11', 'Controller',   'Atomically updates topology.json (.tmp → rename); DU simulators reload within 5 s'),
        ('12', 'Orchestrator', 'JSON-sanitises move result; appends to history; re-calls Gemini [loop continues]'),
        ('13', 'Gemini LLM',   'Repeats move_cell for each remaining overloaded cell until all resolved → loop exits'),
        ('14', 'Gemini LLM',   'Generates summary: cells moved, before/after DU loads'),
        ('15', 'Orchestrator', 'Streams text/plain chunks via sync generator (StreamingResponse)'),
    ])
    prs.save(os.path.join(OUT, 'O1_Query_Flow.pptx'))
    print('Saved O1_Query_Flow.pptx')


# ═══════════════════════════════════════════════════════════════════════════════
# O2 — query_network + move_cell (multi-tool)
# ═══════════════════════════════════════════════════════════════════════════════
def make_o2():
    prs = new_prs()
    title_slide(prs, 'O2',
        'Move the most overloaded cell to a lighter DU.',
        'query_network  →  rank by PRB  →  move_cell  (multi-tool)',
        '4G/5G NSA  ·  Single best move  ·  PRB-based ranking')
    multitool_slide(prs,
        'O2 Query Flow — Single Best Cell Move (Multi-Tool)',
        'fn call 1:\nquery_network\n{}', 'GET /network',
        'fn call 2:\nmove_cell\n(top cell)', 'POST /move/cell',
        ['Rank by prb_dl_pct', '→ top cell + lightest DU', '→ execute move_cell'])
    steps_slide(prs, 'O2 Query — Step-by-Step Trace', [
        ('1',  'User',         'Types query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session_id, message} → Orchestrator /chat  (pure stdlib urllib; --session flag)'),
        ('3',  'Orchestrator', 'Calls build_network_context() → GET /network; appends snapshot to SYSTEM_PROMPT; sends to Gemini'),
        ('4',  'Gemini LLM',   'Tool call 1: query_network()  [while True loop, iteration 1]'),
        ('5',  'Orchestrator', 'GET /network → Controller; yields *[calling tool: query_network...]* inline'),
        ('6',  'Controller',   'Reads topology.json + InfluxDB PRB per cell; returns 30-cell snapshot'),
        ('7',  'Orchestrator', 'JSON-sanitises; appends FunctionResponse to session history; re-calls Gemini [loop iteration 2]'),
        ('8',  'Gemini LLM',   'Ranks cells by prb_dl_pct (desc) → most overloaded cell identified'),
        ('9',  'Gemini LLM',   'Ranks DUs by avg PRB (asc) → lightest target DU identified'),
        ('10', 'Gemini LLM',   'Tool call 2: move_cell(cell_id=X, target_du_id=Y)'),
        ('11', 'Orchestrator', 'POST /move/cell on Controller; yields *[calling tool: move_cell...]* inline'),
        ('12', 'Controller',   'Atomically updates topology.json (.tmp → rename); DU simulators reload in 5 s'),
        ('13', 'Orchestrator', 'JSON-sanitises move result; appends to history; re-calls Gemini → loop exits (no more tools)'),
        ('14', 'Gemini LLM',   'Generates reply: cell moved, from/to DU, expected PRB improvement'),
        ('15', 'Orchestrator', 'Streams text/plain chunks via sync generator (StreamingResponse)'),
    ])
    prs.save(os.path.join(OUT, 'O2_Query_Flow.pptx'))
    print('Saved O2_Query_Flow.pptx')


# ═══════════════════════════════════════════════════════════════════════════════
# I1 — Intent: improve UX near railway station (multi-tool autonomous)
# ═══════════════════════════════════════════════════════════════════════════════
def make_i1():
    prs = new_prs()
    title_slide(prs, 'I1',
        'Improve user experience in the Malleswaram railway station area.',
        'query_network  →  assess area  →  move_cell or plan_network  (autonomous)',
        '4G/5G NSA  ·  Intent-based SON  ·  Gemini decides action autonomously')
    multitool_slide(prs,
        'I1 Query Flow — Intent-Based UX Improvement',
        'fn call 1:\nquery_network\n{}', 'GET /network',
        'fn call 2:\nmove_cell or\nplan_network', 'POST /move/cell\nor /plan',
        ['Assess railway station', 'area cells → decides', 'action autonomously'])
    steps_slide(prs, 'I1 Query — Step-by-Step Trace', [
        ('1',  'User',         'Types intent-based query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session_id, message} → Orchestrator /chat  (pure stdlib urllib; --session flag)'),
        ('3',  'Orchestrator', 'Calls build_network_context() → GET /network; appends snapshot to SYSTEM_PROMPT; sends to Gemini'),
        ('4',  'Gemini LLM',   'Tool call 1: query_network()  [while True loop, iteration 1]'),
        ('5',  'Orchestrator', 'GET /network → Controller; yields *[calling tool: query_network...]* inline'),
        ('6',  'Controller',   'topology.json + InfluxDB KPIs (PRB, SINR, UEs, throughput); returns 30-cell snapshot'),
        ('7',  'Orchestrator', 'JSON-sanitises; appends FunctionResponse to session history; re-calls Gemini [iteration 2]'),
        ('8',  'Gemini LLM',   'Identifies railway station area cells: MLS_RWS_* (lat≈13.008)'),
        ('9',  'Gemini LLM',   'Assesses KPIs: high PRB? low SINR? UEs near max?'),
        ('10', 'Gemini LLM',   'Autonomously decides: move_cell (DU imbalance) or plan_network (capacity gap)'),
        ('11', 'Orchestrator', 'Executes chosen tool; yields *[calling tool...]* inline; updates topology or plan'),
        ('12', 'Orchestrator', 'JSON-sanitises result; appends to history; re-calls Gemini → loop exits'),
        ('13', 'Gemini LLM',   'Generates intent fulfilment summary with expected UX improvements'),
        ('14', 'Orchestrator', 'Streams text/plain chunks via sync generator (StreamingResponse)'),
    ])
    prs.save(os.path.join(OUT, 'I1_Query_Flow.pptx'))
    print('Saved I1_Query_Flow.pptx')


# ═══════════════════════════════════════════════════════════════════════════════
# I2 — Intent: reduce call drops in southern cells (3-tool)
# ═══════════════════════════════════════════════════════════════════════════════
def make_i2():
    prs = new_prs()
    title_slide(prs, 'I2',
        'Reduce call drops in the southern cells.',
        'get_alerts(SINR_LOW)  →  query_network  →  move_cell  (3-tool autonomous)',
        '4G/5G NSA  ·  Intent-based SON  ·  Southern cells: lat < 13.000')
    multitool_slide(prs,
        'I2 Query Flow — Intent-Based Call Drop Reduction',
        'fn call 1:\nget_alerts\n(SINR_LOW)', 'GET /alerts\n?type=SINR_LOW',
        'fn call 2+3:\nquery_network\n+ move_cell', 'GET /network\nPOST /move/cell',
        ['SINR_LOW alerts →', 'confirm sinr_db →', 'rebalance southern cells'])
    steps_slide(prs, 'I2 Query — Step-by-Step Trace', [
        ('1',  'User',         'Types intent-based query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session, message} to Orchestrator /chat'),
        ('3',  'Orchestrator', 'Injects live network snapshot; calls Gemini'),
        ('4',  'Gemini LLM',   'Tool call 1: get_alerts(alert_type="SINR_LOW")'),
        ('5',  'Orchestrator', 'Direct Flux query on InfluxDB: SINR_LOW alerts; yields *[calling tool: get_alerts...]* inline'),
        ('6',  'InfluxDB',     'Returns SINR_LOW records with cell_id, severity, confidence'),
        ('7',  'Orchestrator', 'JSON-sanitises; appends FunctionResponse to session history; re-calls Gemini [iteration 2]'),
        ('8',  'Gemini LLM',   'Tool call 2: query_network() — confirm sinr_db values + cell locations'),
        ('9',  'Orchestrator', 'GET /network → Controller; yields *[calling tool: query_network...]* inline'),
        ('10', 'Orchestrator', 'JSON-sanitises; appends FunctionResponse; re-calls Gemini [iteration 3]'),
        ('11', 'Gemini LLM',   'Filters southern cells (lat < 13.000): MGR, CHD, 6CR, SPG'),
        ('12', 'Gemini LLM',   'Cross-checks SINR_LOW alerts with live sinr_db values'),
        ('13', 'Gemini LLM',   'Tool call 3: move_cell(cell_id=X, target_du_id=Y)'),
        ('14', 'Controller',   'Atomically updates topology.json (.tmp → rename); DU simulators reload in 5 s'),
        ('15', 'Orchestrator', 'JSON-sanitises; appends to history; re-calls Gemini → loop exits; streams reply'),
    ])
    prs.save(os.path.join(OUT, 'I2_Query_Flow.pptx'))
    print('Saved I2_Query_Flow.pptx')


# ═══════════════════════════════════════════════════════════════════════════════
# E1 — Explainability: what actions were taken (no tool)
# ═══════════════════════════════════════════════════════════════════════════════
def make_e1():
    prs = new_prs()
    title_slide(prs, 'E1',
        'What actions did you take to optimize the network?',
        'No tool call  →  LLM summarises from conversation history',
        'Context: follow-up after O1/O2 optimization session')
    no_tool_slide(prs, 'E1 Query Flow — Action Summary (No Tool)',
        ['Reads all tool call', 'records from history', '→ ordered action list'])
    steps_slide(prs, 'E1 Query — Step-by-Step Trace', [
        ('1',  'User',         'Types explainability query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session_id, message} → Orchestrator /chat  (pure stdlib urllib; --session flag for named sessions)'),
        ('3',  'Orchestrator', 'Calls build_network_context() → GET /network on Controller; appends live cell snapshot to SYSTEM_PROMPT'),
        ('4',  'Orchestrator', 'Sends SYSTEM_PROMPT (static) + live snapshot (dynamic) + session history (types.Content) to Gemini'),
        ('5',  'Gemini LLM',   'No tool call — reads all prior tool call records + results from session history (types.Content list)'),
        ('6',  'Gemini LLM',   'Extracts: which tools called, with what args, and outcomes'),
        ('7',  'Gemini LLM',   'Reconstructs: alerts detected → cells moved → topology updated'),
        ('8',  'Gemini LLM',   'Generates ordered action summary with before/after metrics; loop exits (no tool calls)'),
        ('9',  'Orchestrator', 'Streams text/plain chunks via sync generator (Starlette thread pool → StreamingResponse)'),
        ('10', 'chat.py',      'Prints streamed action summary to operator terminal'),
    ])
    prs.save(os.path.join(OUT, 'E1_Query_Flow.pptx'))
    print('Saved E1_Query_Flow.pptx')


# ═══════════════════════════════════════════════════════════════════════════════
# E2 — Explainability: why was that cell moved (no tool)
# ═══════════════════════════════════════════════════════════════════════════════
def make_e2():
    prs = new_prs()
    title_slide(prs, 'E2',
        'Why did you move that cell to a different DU?',
        'No tool call  →  LLM explains specific move decision from context',
        'Context: follow-up after a move_cell action')
    no_tool_slide(prs, 'E2 Query Flow — Move Decision Explanation (No Tool)',
        ['Reads move_cell context:', 'alert + DU loads at time', '→ causal explanation'])
    steps_slide(prs, 'E2 Query — Step-by-Step Trace', [
        ('1',  'User',         'Types explainability query into chat.py'),
        ('2',  'chat.py',      'HTTP POST {session_id, message} → Orchestrator /chat  (pure stdlib urllib; --session flag for named sessions)'),
        ('3',  'Orchestrator', 'Calls build_network_context() → GET /network on Controller; appends live cell snapshot to SYSTEM_PROMPT'),
        ('4',  'Orchestrator', 'Sends SYSTEM_PROMPT (static) + live snapshot (dynamic) + session history (types.Content) to Gemini'),
        ('5',  'Gemini LLM',   'No tool call — reconstructs decision from move_cell record in session history (types.Content list)'),
        ('6',  'Gemini LLM',   'Identifies: which cell, source DU, target DU, triggering alert'),
        ('7',  'Gemini LLM',   'Recalls DU load comparison at time of decision'),
        ('8',  'Gemini LLM',   'Generates causal explanation: overload alert → DU ranking → move; loop exits'),
        ('9',  'Orchestrator', 'Streams text/plain chunks via sync generator (Starlette thread pool → StreamingResponse)'),
        ('10', 'chat.py',      'Prints streamed explanation to operator terminal'),
    ])
    prs.save(os.path.join(OUT, 'E2_Query_Flow.pptx'))
    print('Saved E2_Query_Flow.pptx')


# ─── Run all ──────────────────────────────────────────────────────────────────
if __name__ == '__main__':
    make_f1(); make_f2()
    make_a1(); make_a2()
    make_o1(); make_o2()
    make_i1(); make_i2()
    make_e1(); make_e2()
    print('\nAll 10 PPTs generated.')
